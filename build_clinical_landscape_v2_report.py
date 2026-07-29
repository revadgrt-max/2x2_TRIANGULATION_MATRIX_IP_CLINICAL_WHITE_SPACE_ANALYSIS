"""
Build the PDF + PPTX deliverables for the Clinical Landscape V2 deep-dive.

Parses the already-executed notebook (input/clinical_landscape_final_version2.ipynb)
directly -- no recomputation -- so every number, table and chart in the deliverables
mirrors exactly what is in the notebook. Every table cell output was saved by
build_clinical_landscape_v2_notebook.py with its clean name in `output.metadata.name`,
matching the CSV already sitting in output/clinical2_landscape_final_version_2/data/ --
so this script re-uses those exact files rather than re-exporting under new names.

Produces, under output/clinical2_landscape_final_version_2/:
    Clinical_Landscape_V2_Report.pdf   - full report: title, clickable Table of Contents,
                                         TL;DR, Executive Summary (incl. post-interview V2.1
                                         addendum), then one section per notebook Entry
                                         (SQL/description note + charts + first-20-rows
                                         tables with a link to the full CSV in data/)
    Clinical_Landscape_V2_Deck.pptx    - slide deck: title, clickable agenda/TOC, TL;DR,
                                         executive summary, then one slide per Entry
                                         (chart or table snapshot)
    data/*.csv                        - already present (one CSV per table, all rows)
    figures/*.png                     - already present (one PNG per chart)
"""
import json
import re
from datetime import date
from pathlib import Path

import pandas as pd

from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table as RLTable,
    TableStyle, PageBreak,
)
from PIL import Image as PILImage

# ============================================================================
# Paths
# ============================================================================
BASE = Path("/Users/revathisekar/Documents/vibe analytics trial 2")
NOTEBOOK_SRC = BASE / "input" / "clinical_landscape_final_version2.ipynb"
OUT = BASE / "output" / "clinical2_landscape_final_version_2"
DATA = OUT / "data"
FIG = OUT / "figures"
DATA.mkdir(parents=True, exist_ok=True)
FIG.mkdir(parents=True, exist_ok=True)

TODAY = date(2026, 7, 27).strftime("%B %d, %Y")
REPORT_TITLE = "Clinical Trial Landscape V2 -- Antibody Therapeutics in Oncology"
SUBTITLE = "Deep-Dive: Modality x Target x Tumor-Type x Phase, Subtle Divergences, Whitespace & Risk-Tiered Portfolio Recommendation"


def pd_isna(v):
    try:
        return pd.isna(v)
    except (TypeError, ValueError):
        return False


# ============================================================================
# 1. Parse the notebook into ordered "Entry" sections
# ============================================================================
print("Loading and parsing notebook ...")
nb = json.loads(NOTEBOOK_SRC.read_text())


def cell_source(cell):
    src = cell.get("source", "")
    return "".join(src) if isinstance(src, list) else src


ENTRY_RE = re.compile(r"^##\s+Entry\s+(\d+)\s*:\s*(.+)$")

entries = []
current = None
intro_text = None
_fig_counter = 0

for cell in nb["cells"]:
    ctype = cell.get("cell_type")
    src = cell_source(cell)
    if ctype == "markdown":
        stripped = src.strip()
        lines = stripped.split("\n")
        first_line = lines[0].strip() if lines else ""
        m = ENTRY_RE.match(first_line)
        if m:
            if current is not None:
                entries.append(current)
            number = int(m.group(1))
            title = m.group(2).strip()
            desc = "\n".join(lines[1:]).strip()
            current = {"number": number, "title": title, "description": desc, "items": []}
        elif current is None and intro_text is None:
            intro_text = stripped
        continue
    elif ctype == "code":
        if current is None:
            continue
        for out in cell.get("outputs", []):
            otype = out.get("output_type")
            data = out.get("data", {})
            meta = out.get("metadata", {}) or {}
            if otype in ("display_data", "execute_result"):
                if "image/png" in data:
                    b64 = data["image/png"]
                    if isinstance(b64, list):
                        b64 = "".join(b64)
                    import base64
                    _fig_counter += 1
                    fname = f"entry{current['number']:02d}_fig{_fig_counter:02d}.png"
                    fpath = FIG / fname
                    if not fpath.exists():
                        fpath.write_bytes(base64.b64decode(b64))
                    current["items"].append({"type": "image", "path": fpath})
                    continue
                html = data.get("text/html")
                if html:
                    html = "".join(html) if isinstance(html, list) else html
                    try:
                        import io as _io
                        dfs = pd.read_html(_io.StringIO(html))
                        df = dfs[0] if dfs else None
                    except Exception:
                        df = None
                    if df is not None and len(df.columns) > 0:
                        tname = meta.get("name") or f"entry{current['number']:02d}_table{len([i for i in current['items'] if i['type']=='table'])+1:02d}"
                        safe = re.sub(r"[^A-Za-z0-9_]+", "_", tname)
                        csv_name = f"{safe}.csv"
                        csv_path = DATA / csv_name
                        if not csv_path.exists():
                            df.to_csv(csv_path, index=False)
                        current["items"].append({"type": "table", "df": df, "csv_name": csv_name, "title": tname})
                        continue
                txt = data.get("text/plain")
                if txt:
                    txt = "".join(txt) if isinstance(txt, list) else txt
                    if txt.strip():
                        current["items"].append({"type": "text", "text": txt.strip()})
                        continue
            elif otype == "stream":
                txt_raw = out.get("text", "")
                txt = "".join(txt_raw) if isinstance(txt_raw, list) else txt_raw
                if txt.strip():
                    current["items"].append({"type": "text", "text": txt.strip()})

if current is not None:
    entries.append(current)

print(f"Parsed {len(entries)} entries from notebook.")
n_figs = sum(1 for e in entries for it in e["items"] if it["type"] == "image")
n_tables = sum(1 for e in entries for it in e["items"] if it["type"] == "table")
print(f"Extracted {n_figs} figures and {n_tables} tables (CSVs already on disk in data/, reused where present).")

# ============================================================================
# 2. TL;DR + Executive Summary text (pulled verbatim from the notebook's own
#    printed summaries -- Entry 36 "Executive Summary" and Entry 43 "Final
#    Consolidated Recommendation Matrix & Updated Executive Summary V2.1")
# ============================================================================
e36 = next((e for e in entries if e["number"] == 36), None)
e43 = next((e for e in entries if e["number"] == 43), None)
exec_summary_1 = next((it["text"] for it in (e36["items"] if e36 else []) if it["type"] == "text" and "EXECUTIVE SUMMARY" in it["text"]), "")
exec_summary_2_full = next((it["text"] for it in (e43["items"] if e43 else []) if it["type"] == "text" and "EXECUTIVE SUMMARY V2.1" in it["text"]), "")
exec_summary_2 = exec_summary_2_full.split("EXECUTIVE SUMMARY V2.1")[-1].strip() if exec_summary_2_full else ""
if exec_summary_2:
    exec_summary_2 = "EXECUTIVE SUMMARY V2.1" + exec_summary_2 if not exec_summary_2.startswith("--") else "EXECUTIVE SUMMARY V2.1 " + exec_summary_2


def grab(pattern, text, cast=str, default=None):
    m = re.search(pattern, text)
    if not m:
        return default
    val = m.group(1).replace(",", "")
    try:
        return cast(val)
    except Exception:
        return default


N_TRIALS = grab(r"Universe analyzed:\s*([\d,]+)\s*in-scope", exec_summary_1, int)
N_TARGETS_5 = grab(r"Among\s*(\d+)\s*targets", exec_summary_1, int)
N_VALIDATED = grab(r"(\d+)\s*are VALIDATED", exec_summary_1, int)
N_MIXED = grab(r"(\d+)\s*are MIXED", exec_summary_1, int)
N_FAILED_ONLY = grab(r"(\d+)\s*have FAILED ONLY", exec_summary_1, int)
N_STILL = grab(r"(\d+)\s*are STILL IN TRIAL", exec_summary_1, int)
N_REGIMEN_FLAG = grab(r"(\d+)\s*target x tumor-type pairs show a population-matched IO-combo", exec_summary_1, int)
N_MOD_FLAG = grab(r"(\d+)\s*target x tumor-type pairs show a population-matched naked-mAb", exec_summary_1, int)
N_HEATING = grab(r"(\d+)\s*targets are 'heating up'", exec_summary_1, int)
N_WHITESPACE = grab(r"(\d+)\s*targets qualify as Phase-1/2-only", exec_summary_1, int)
TOP_WS_TARGET = grab(r"Top-ranked:\s*([A-Za-z0-9\|\-\.]+)\.", exec_summary_1, str)
N_SPOTLIGHT = grab(r"(\d+)\s*pass the tighter NOVEL-CANDIDATE", exec_summary_1, int)

N_LOW_RISK = grab(r"(\d+)\s*low-risk", exec_summary_2, int)
N_MED_RISK = grab(r"(\d+)\s*medium-risk", exec_summary_2, int)
N_HIGH_RISK = grab(r"(\d+)\s*high-risk", exec_summary_2, int)
N_MONO_READY = grab(r"(\d+)\s*candidates\s*across both views are monotherapy-ready", exec_summary_2, int)
N_IO_NEEDED = grab(r"(\d+)\s*would need IO-combo access", exec_summary_2, int)

TLDR_BULLETS = [
    f"{N_TRIALS:,} in-scope antibody-modality oncology trial records analyzed (mAb, ADC, bispecific, BiTE, CAR-T, radioligand, PROTAC), standardized on `target_harmonized`." if N_TRIALS else
    "Full antibody-modality oncology trial universe analyzed, standardized on `target_harmonized`.",
    f"Of {N_TARGETS_5} targets with >=5 trials: {N_VALIDATED} validated, {N_MIXED} mixed (win in some settings, lose in others), {N_FAILED_ONLY} failed-only, {N_STILL} still awaiting readout." if N_TARGETS_5 else
    "Targets classified into validated / mixed / failed-only / still-in-trial status.",
    f"{N_REGIMEN_FLAG} target x tumor-type pairs show a population-matched IO-combo vs all-comers divergence; {N_MOD_FLAG} show a naked-mAb vs ADC divergence -- the SAME target, SAME tumor type, opposite outcome depending on regimen/modality." if N_REGIMEN_FLAG is not None else
    "Population-matched divergence analysis flags targets whose outcome flips by regimen or modality within the same tumor type.",
    f"{N_HEATING} targets are 'heating up' (accelerating recent activity or a brand-new wave since 2022)." if N_HEATING else
    "Momentum analysis identifies targets with accelerating recent trial activity.",
    f"{N_WHITESPACE} targets qualify as Phase-1/2-only clinical whitespace (zero recorded failures); top-ranked: {TOP_WS_TARGET}. {N_SPOTLIGHT} pass the tighter novel-candidate filter (first-in-human since 2020, <=3 sponsors)." if N_WHITESPACE else
    "Whitespace scoring surfaces Phase-1/2-only targets with zero recorded failures.",
    f"Post-interview: risk-tiered into {N_LOW_RISK} low-risk (validated biology), {N_MED_RISK} medium-risk (hybrid: validated target + novel modality angle -- ADC-first, matching in-house capability), {N_HIGH_RISK} high-risk (whitespace/novel) candidates." if N_LOW_RISK is not None else
    "Post-interview addendum adds a Low/Medium/High risk-tier framework.",
    f"Combo-strategy feasibility tagged per candidate: {N_MONO_READY} are monotherapy-ready today; {N_IO_NEEDED} would require IO-combo access to realize their only recorded success." if N_MONO_READY is not None else
    "Every shortlisted candidate is tagged with its combo-strategy feasibility (monotherapy-ready vs requires IO/chemo/TKI-combo access).",
]

print("TL;DR bullets prepared:")
for b in TLDR_BULLETS:
    print(" -", b[:100])

# ============================================================================
# 3. PDF report -- clickable internal Table of Contents
# ============================================================================
print("Building PDF report ...")

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(name="TitleBig", fontSize=22, leading=26, alignment=1, textColor=colors.HexColor("#1F2A44"), spaceAfter=10))
styles.add(ParagraphStyle(name="SubtitleBig", fontSize=12, leading=15, alignment=1, textColor=colors.HexColor("#444444"), spaceAfter=6, fontName="Helvetica-Oblique"))
styles.add(ParagraphStyle(name="MetaCenter", fontSize=9.5, leading=13, alignment=1, textColor=colors.HexColor("#666666")))
styles.add(ParagraphStyle(name="H1", fontSize=16, leading=20, spaceBefore=14, spaceAfter=8, textColor=colors.HexColor("#1F2A44")))
styles.add(ParagraphStyle(name="H2", fontSize=12.5, leading=16, spaceBefore=10, spaceAfter=6, textColor=colors.HexColor("#2563EB")))
styles.add(ParagraphStyle(name="Body9", fontSize=9, leading=12.5, spaceAfter=6))
styles.add(ParagraphStyle(name="Note9", fontSize=9, leading=12.5, spaceAfter=6, textColor=colors.HexColor("#333333"), fontName="Helvetica-Oblique"))
styles.add(ParagraphStyle(name="Mono8", fontSize=7.5, leading=10.5, spaceAfter=6, fontName="Courier"))
styles.add(ParagraphStyle(name="Bullet9", fontSize=9, leading=12.5, leftIndent=14, bulletIndent=2, spaceAfter=4))
styles.add(ParagraphStyle(name="TOCLink", fontSize=10, leading=15, textColor=colors.HexColor("#1155CC")))


def esc(text):
    return str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def esc_pre(text):
    return esc(text).replace("\n", "<br/>")


def pdf_df_table(df, max_rows=20, font_size=6.3):
    shown = df.head(max_rows)
    data = [[esc(c) for c in shown.columns]] + [
        ["" if pd_isna(v) else esc(v) for v in row] for row in shown.itertuples(index=False)
    ]
    ncols = len(shown.columns)
    avail_width = 7.3 * inch
    col_width = avail_width / max(ncols, 1)
    t = RLTable(data, colWidths=[col_width] * ncols, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F2A44")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, -1), font_size),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CCCCCC")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F6FA")]),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 3),
        ("RIGHTPADDING", (0, 0), (-1, -1), 3),
    ]))
    return t


def pdf_render_items(story, items):
    for it in items:
        if it["type"] == "text":
            story.append(Paragraph(esc_pre(it["text"]), styles["Mono8"] if len(it["text"]) > 300 else styles["Note9"]))
            story.append(Spacer(1, 4))
        elif it["type"] == "table":
            story.append(pdf_df_table(it["df"], max_rows=20))
            rel = f"data/{it['csv_name']}"
            if len(it["df"]) > 20:
                story.append(Paragraph(
                    f"Showing first 20 of {len(it['df']):,} rows. "
                    f"<link href=\"{rel}\">Click to open full CSV ({esc(it['csv_name'])})</link>",
                    styles["Body9"]))
            else:
                story.append(Paragraph(
                    f"All {len(it['df']):,} rows shown above. <link href=\"{rel}\">Click to open CSV ({esc(it['csv_name'])})</link>",
                    styles["Body9"]))
            story.append(Spacer(1, 8))
        elif it["type"] == "image":
            with PILImage.open(it["path"]) as im:
                w, h = im.size
            max_w = 6.6 * inch
            ratio = h / w
            draw_w = max_w
            draw_h = max_w * ratio
            max_h = 7.3 * inch
            if draw_h > max_h:
                draw_h = max_h
                draw_w = max_h / ratio
            story.append(RLImage(str(it["path"]), width=draw_w, height=draw_h))
            story.append(Spacer(1, 8))


toc_entries = [("exec_summary", "Executive Summary"), ("tldr", "TL;DR")]
toc_entries += [(f"sec_entry_{e['number']}", f"Entry {e['number']}: {e['title']}") for e in entries]
toc_entries += [("appendix", "Appendix: Data Files Index")]

story = []
story.append(Spacer(1, 1.3 * inch))
story.append(Paragraph(esc(REPORT_TITLE), styles["TitleBig"]))
story.append(Paragraph(esc(SUBTITLE), styles["SubtitleBig"]))
story.append(Spacer(1, 0.25 * inch))
story.append(Paragraph(
    f"Report generated: {TODAY}<br/>Source notebook: clinical_landscape_final_version2.ipynb<br/>"
    f"Source data: clinical_final_version1.csv" + (f" ({N_TRIALS:,} in-scope trial records)" if N_TRIALS else ""),
    styles["MetaCenter"],
))
story.append(PageBreak())

story.append(Paragraph('<a name="toc"/>Table of Contents', styles["H1"]))
for anchor, label in toc_entries:
    story.append(Paragraph(f'<link href="#{anchor}">{esc(label)}</link>', styles["TOCLink"]))
story.append(PageBreak())

story.append(Paragraph('<a name="exec_summary"/>Executive Summary', styles["H1"]))
if exec_summary_1:
    story.append(Paragraph(esc_pre(exec_summary_1), styles["Mono8"]))
    story.append(Spacer(1, 8))
if exec_summary_2:
    story.append(Paragraph(esc_pre(exec_summary_2), styles["Mono8"]))
story.append(PageBreak())

story.append(Paragraph('<a name="tldr"/>TL;DR', styles["H1"]))
for b in TLDR_BULLETS:
    story.append(Paragraph(f"&bull; {esc(b)}", styles["Bullet9"]))
story.append(PageBreak())

for e in entries:
    story.append(Paragraph(f'<a name="sec_entry_{e["number"]}"/>Entry {e["number"]}: {esc(e["title"])}', styles["H1"]))
    if e["description"]:
        story.append(Paragraph(esc(e["description"]), styles["Body9"]))
    pdf_render_items(story, e["items"])
    story.append(PageBreak())

story.append(Paragraph('<a name="appendix"/>Appendix: Data Files Index', styles["H1"]))
story.append(Paragraph(esc(
    "Every table in this report -- regardless of size -- has its full underlying data deposited as a "
    "clearly labelled CSV file in the data/ folder next to this report."
), styles["Body9"]))
appendix_rows = [["Entry", "Rows", "File"]]
for e in entries:
    for it in e["items"]:
        if it["type"] == "table":
            appendix_rows.append([f"Entry {e['number']}: {e['title'][:40]}", f"{len(it['df']):,}", it["csv_name"]])
t = RLTable(appendix_rows, colWidths=[3.2 * inch, 0.8 * inch, 3.3 * inch], repeatRows=1)
t.setStyle(TableStyle([
    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F2A44")),
    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
    ("FONTSIZE", (0, 0), (-1, -1), 7),
    ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CCCCCC")),
    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F6FA")]),
]))
story.append(t)

pdf_path = OUT / "Clinical_Landscape_V2_Report.pdf"
pdf_doc = SimpleDocTemplate(str(pdf_path), pagesize=letter, topMargin=0.6 * inch, bottomMargin=0.6 * inch,
                            leftMargin=0.6 * inch, rightMargin=0.6 * inch, title=REPORT_TITLE)
pdf_doc.build(story)
print(f"Saved PDF: {pdf_path}")

# ============================================================================
# 4. PPTX -- 16:9 deck with a clickable agenda (internal hyperlinks to slides)
# ============================================================================
print("Building PPTX ...")

prs = Presentation()
prs.slide_width = PIn(13.333)
prs.slide_height = PIn(7.5)
BLANK = prs.slide_layouts[6]
DARK = PRGBColor(0x1F, 0x2A, 0x44)
ACCENT = PRGBColor(0x25, 0x63, 0xEB)
GREY = PRGBColor(0x44, 0x44, 0x44)


def add_title_box(slide, text, top=PIn(0.35), size=30, color=DARK, height=PIn(1.0)):
    box = slide.shapes.add_textbox(PIn(0.5), top, prs.slide_width - PIn(1.0), height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PPt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return box


def add_bullets(slide, bullets, top=PIn(1.5), left=PIn(0.7), width=None, height=None, size=16):
    width = width or (prs.slide_width - PIn(1.4))
    height = height or (prs.slide_height - top - PIn(0.4))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022 {b}"
        p.font.size = PPt(size)
        p.font.color.rgb = GREY
        p.space_after = PPt(10)
    return box


def add_image_slide(title, image_path, caption=None):
    slide = prs.slides.add_slide(BLANK)
    add_title_box(slide, title, size=22, height=PIn(0.7))
    with PILImage.open(image_path) as im:
        w, h = im.size
    ratio = h / w
    max_w = PIn(12.2)
    max_h = PIn(5.9)
    draw_w = max_w
    draw_h = PIn(max_w.inches * ratio)
    if draw_h > max_h:
        draw_h = max_h
        draw_w = PIn(max_h.inches / ratio)
    left = (prs.slide_width - draw_w) / 2
    slide.shapes.add_picture(str(image_path), left, PIn(1.15), width=draw_w, height=draw_h)
    if caption:
        cap = slide.shapes.add_textbox(PIn(0.5), prs.slide_height - PIn(0.55), prs.slide_width - PIn(1.0), PIn(0.4))
        p = cap.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = PPt(11)
        p.font.italic = True
        p.font.color.rgb = GREY
        p.alignment = PP_ALIGN.CENTER
    return slide


def add_table_slide(title, df, max_rows=8):
    slide = prs.slides.add_slide(BLANK)
    add_title_box(slide, title, size=22, height=PIn(0.7))
    shown = df.head(max_rows)
    rows, cols = shown.shape[0] + 1, shown.shape[1]
    left, top, width, height = PIn(0.6), PIn(1.3), prs.slide_width - PIn(1.2), PIn(5.6)
    gtable = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for j, col in enumerate(shown.columns):
        gtable.cell(0, j).text = str(col)
    for i, (_, row) in enumerate(shown.iterrows(), start=1):
        for j, val in enumerate(row):
            gtable.cell(i, j).text = "" if pd_isna(val) else str(val)
    return slide


def add_text_slide(title, text, size=13):
    slide = prs.slides.add_slide(BLANK)
    add_title_box(slide, title, size=26)
    lines = [ln for ln in text.split("\n") if ln.strip()][:22]
    add_bullets(slide, lines, size=size, height=PIn(5.8))
    return slide


# Title slide
slide = prs.slides.add_slide(BLANK)
box = slide.shapes.add_textbox(PIn(0.8), PIn(2.4), prs.slide_width - PIn(1.6), PIn(1.6))
p = box.text_frame.paragraphs[0]
p.text = REPORT_TITLE
p.font.size = PPt(34)
p.font.bold = True
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER
sub_box = slide.shapes.add_textbox(PIn(0.8), PIn(3.7), prs.slide_width - PIn(1.6), PIn(1.4))
p2 = sub_box.text_frame.paragraphs[0]
p2.text = SUBTITLE
p2.font.size = PPt(16)
p2.font.italic = True
p2.font.color.rgb = GREY
p2.alignment = PP_ALIGN.CENTER
meta_box = slide.shapes.add_textbox(PIn(0.8), PIn(5.1), prs.slide_width - PIn(1.6), PIn(0.8))
p3 = meta_box.text_frame.paragraphs[0]
p3.text = f"Generated {TODAY}  |  clinical_landscape_final_version2.ipynb  |  clinical_final_version1.csv"
p3.font.size = PPt(13)
p3.font.color.rgb = GREY
p3.alignment = PP_ALIGN.CENTER

# TL;DR slide
add_text_slide("TL;DR", "\n".join(TLDR_BULLETS), size=14)

# Executive summary slide(s) -- split long printed text across slides of ~18 lines
for label, text in (("Executive Summary", exec_summary_1), ("Executive Summary V2.1 (Post-Interview)", exec_summary_2)):
    if not text:
        continue
    lines = [ln for ln in text.split("\n") if ln.strip()]
    chunk_size = 16
    for i in range(0, len(lines), chunk_size):
        chunk = lines[i:i + chunk_size]
        title = label if i == 0 else f"{label} (cont.)"
        add_text_slide(title, "\n".join(chunk), size=13)

# Agenda / TOC slide -- clickable, jumps to each Entry's first slide
agenda_slide = prs.slides.add_slide(BLANK)
add_title_box(agenda_slide, "Agenda", size=32)
agenda_top = PIn(1.4)
agenda_left_col = [PIn(0.6), PIn(6.9)]
col_width = PIn(6.0)
row_height = PIn(0.34)
n_per_col = 22

# Build entry slides first (so we know slide references), collecting the first slide of each entry.
entry_first_slide = {}
entry_slides_data = []
for e in entries:
    images = [it for it in e["items"] if it["type"] == "image"]
    tables = [it for it in e["items"] if it["type"] == "table"]
    title = f"Entry {e['number']}: {e['title']}"
    made_slides = []
    if images:
        for k, img in enumerate(images):
            cap = title if k == 0 else f"{title} (cont.)"
            made_slides.append(add_image_slide(title, img["path"], caption=e["title"]))
    elif tables:
        made_slides.append(add_table_slide(title, tables[0]["df"], max_rows=8))
    else:
        notes = [it["text"] for it in e["items"] if it["type"] == "text"]
        text_blob = "\n".join(notes) if notes else "See full notebook / PDF report for details."
        made_slides.append(add_text_slide(title, text_blob, size=13))
    if made_slides:
        entry_first_slide[e["number"]] = made_slides[0]

# Now populate the agenda slide with per-entry clickable textboxes (one shape per
# line so each line's click_action can target its own slide -- python-pptx only
# supports one click_action per shape).
for idx, e in enumerate(entries):
    col = idx // n_per_col
    row = idx % n_per_col
    left = agenda_left_col[min(col, len(agenda_left_col) - 1)]
    top = agenda_top + PIn(row_height.inches * row)
    box = agenda_slide.shapes.add_textbox(left, top, col_width, row_height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Entry {e['number']}: {e['title']}"
    p.font.size = PPt(11)
    p.font.color.rgb = ACCENT
    target = entry_first_slide.get(e["number"])
    if target is not None:
        box.click_action.target_slide = target

pptx_path = OUT / "Clinical_Landscape_V2_Deck.pptx"
prs.save(pptx_path)
print(f"Saved PPTX: {pptx_path}")

print("\nDone.")
print(f"PDF:  {pdf_path}")
print(f"PPTX: {pptx_path}")
print(f"Data: {DATA} ({len(list(DATA.glob('*.csv')))} CSVs)")
print(f"Figures: {FIG} ({len(list(FIG.glob('*.png')))} PNGs)")
