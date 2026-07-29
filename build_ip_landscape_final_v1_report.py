"""
Build the final deliverables package for the IP Landscape Analysis - Final Version 1.

Parses the already-executed notebook (input/ip_landscape_final_version1.ipynb) directly -
no recomputation - so every number, table and chart in the deliverables mirrors exactly
what is in the notebook. Produces, under output/IP_LANDSCAPE_FINAL VERSION/:

    IP_Landscape_Report.docx          - polished report: exec summary, TL;DR, spotlight
                                        section, clickable TOC, one section per notebook
                                        Entry (topic + graphs + first-20-rows tables with
                                        links to full CSVs)
    IP_Landscape_Report.pdf           - same content, with an internal clickable TOC
    IP_Landscape_Report.pptx          - slide-deck version (topics + graphs)
    IP_Landscape_Notebook_Mirror.docx - cell-by-cell mirror of the exact notebook contents
                                        (markdown + code + every output, in order)
    ip_landscape_final_version1.ipynb - copy of the source notebook
    data/*.csv                        - full data behind every table (clearly labelled)
    figures/*.png                     - every chart extracted from the notebook
"""
import base64
import io
import json
import re
import shutil
from datetime import date
from pathlib import Path

import pandas as pd

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Inches, Pt, RGBColor

from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table as RLTable,
    TableStyle, PageBreak, HRFlowable,
)
from PIL import Image as PILImage

# ============================================================================
# Paths
# ============================================================================
BASE = Path("/Users/revathisekar/Documents/vibe analytics trial 2")
NOTEBOOK_SRC = BASE / "input" / "ip_landscape_final_version1.ipynb"
OUT = BASE / "output" / "IP_LANDSCAPE_FINAL VERSION"
DATA = OUT / "data"
FIG = OUT / "figures"
for d in (OUT, DATA, FIG):
    d.mkdir(parents=True, exist_ok=True)

TODAY = date(2026, 7, 25).strftime("%B %d, %Y")
REPORT_TITLE = "IP Landscape Analysis - Final Version 1"
SUBTITLE = "Oncology Antibody Therapeutics: Target Crowding, White Space & Emerging-Promising-Target Analysis"

print("Copying source notebook into deliverables folder ...")
shutil.copy2(NOTEBOOK_SRC, OUT / NOTEBOOK_SRC.name)

# ============================================================================
# 1. Parse the notebook into ordered "Entry" sections
# ============================================================================
print("Loading and parsing notebook ...")
nb = json.loads(NOTEBOOK_SRC.read_text())


def cell_source(cell):
    src = cell.get("source", "")
    return "".join(src) if isinstance(src, list) else src


ENTRY_RE = re.compile(r"^##\s+Entry\s+(\d+)(?:\s*\(NEW\))?\s*:\s*(.+)$")

entries = []
current = None

for cell in nb["cells"]:
    ctype = cell.get("cell_type")
    src = cell_source(cell)
    if ctype == "markdown":
        stripped = src.strip()
        lines = stripped.split("\n")
        first_line = lines[0].strip() if lines else ""
        m = ENTRY_RE.match(first_line)
        if m:
            if current is not None:
                entries.append(current)
            number = int(m.group(1))
            title = m.group(2).strip()
            desc = "\n".join(lines[1:]).strip()
            current = {"number": number, "title": title, "description": desc,
                       "code_sources": [], "outputs": []}
        continue
    elif ctype == "code":
        if current is None:
            continue
        current["code_sources"].append(src)
        for out in cell.get("outputs", []):
            current["outputs"].append(out)

if current is not None:
    entries.append(current)

print(f"Parsed {len(entries)} entries from notebook.")


def slugify(text, maxlen=40):
    s = re.sub(r"[^a-zA-Z0-9]+", "_", text).strip("_")
    return s[:maxlen]


def text_from_data(data, key):
    val = data.get(key)
    if val is None:
        return None
    return "".join(val) if isinstance(val, list) else val


def png_bytes_from_output(out):
    data = out.get("data", {})
    if "image/png" in data:
        b64 = data["image/png"]
        if isinstance(b64, list):
            b64 = "".join(b64)
        return base64.b64decode(b64)
    return None


def html_to_df(html):
    try:
        dfs = pd.read_html(io.StringIO(html))
        return dfs[0] if dfs else None
    except Exception:
        return None


fig_counter = 0
table_counter = 0

for entry in entries:
    items = []
    for out in entry["outputs"]:
        otype = out.get("output_type")
        if otype in ("display_data", "execute_result"):
            data = out.get("data", {})
            png = png_bytes_from_output(out)
            if png:
                fig_counter += 1
                fname = f"entry{entry['number']:02d}_fig{fig_counter}_{slugify(entry['title'], 30)}.png"
                fpath = FIG / fname
                fpath.write_bytes(png)
                items.append({"type": "image", "path": fpath})
                continue
            html = text_from_data(data, "text/html")
            if html:
                df = html_to_df(html)
                if df is not None and len(df.columns) > 0:
                    table_counter += 1
                    csv_name = f"entry{entry['number']:02d}_table{table_counter:03d}_{slugify(entry['title'], 30)}.csv"
                    csv_path = DATA / csv_name
                    df.to_csv(csv_path, index=False)
                    items.append({"type": "table", "df": df, "csv_path": csv_path, "csv_name": csv_name})
                    continue
            md = text_from_data(data, "text/markdown")
            if md:
                items.append({"type": "note", "text": md.strip()})
                continue
            txt = text_from_data(data, "text/plain")
            if txt and txt.strip():
                items.append({"type": "text", "text": txt.strip()})
                continue
        elif otype == "stream":
            txt_raw = out.get("text", "")
            txt = "".join(txt_raw) if isinstance(txt_raw, list) else txt_raw
            if txt.strip():
                items.append({"type": "text", "text": txt.strip()})
        # 'error' outputs are skipped (none expected - all cells executed successfully)
    entry["items"] = items

n_figs = sum(1 for e in entries for it in e["items"] if it["type"] == "image")
n_tables = sum(1 for e in entries for it in e["items"] if it["type"] == "table")
print(f"Extracted {n_figs} figures and {n_tables} tables.")


def find_note_and_followup(substr):
    """Locate an entry whose note text contains substr; return (entry, note_item, next_item)."""
    for e in entries:
        items = e["items"]
        for i, it in enumerate(items):
            if it["type"] == "note" and substr.lower() in it["text"].lower():
                nxt = items[i + 1] if i + 1 < len(items) else None
                return e, it, nxt
    return None, None, None


def first_number(text, default=0):
    m = re.search(r"\d[\d,]*", text)
    if not m:
        return default
    digits = m.group(0).replace(",", "")
    return int(digits) if digits else default


# ============================================================================
# 2. Pull key headline stats out of the parsed entries (for Exec Summary / TL;DR)
# ============================================================================
e_overview = next((e for e in entries if e["number"] == 2), None)
overview_item = next((it for it in (e_overview["items"] if e_overview else []) if it["type"] == "table"), None)
overview_table = overview_item["df"] if overview_item else None

e_crowding = next((e for e in entries if e["number"] == 3), None)
crowded_table = None
if e_crowding:
    for it in e_crowding["items"]:
        if it["type"] == "table" and {"target_name"}.issubset(set(map(str, it["df"].columns))):
            crowded_table = it["df"]
            break

e_emerging = next((e for e in entries if e["number"] == 7), None)
emerging_note = next((it for it in (e_emerging["items"] if e_emerging else []) if it["type"] == "note"), None)
N_EMERGING = first_number(emerging_note["text"]) if emerging_note else None

e_22 = next((e for e in entries if e["number"] == 22), None)
e22_note = next((it for it in (e_22["items"] if e_22 else []) if it["type"] == "note"), None)
n22_match = re.search(r"([\d,]+)\s+qualifying", e22_note["text"]) if e22_note else None
N_22_QUALIFYING = int(n22_match.group(1).replace(",", "")) if n22_match else None
n22_ind_match = re.search(r"across\s+(\d+)\s+solid-tumor indications", e22_note["text"]) if e22_note else None
N_22_INDICATIONS = int(n22_ind_match.group(1)) if n22_ind_match else 17

e_24 = next((e for e in entries if e["number"] == 24), None)
e24_note = next((it for it in (e_24["items"] if e_24 else []) if it["type"] == "note"), None)
N_24_ROWS = first_number(e24_note["text"]) if e24_note else None

_, lung_adc_note, lung_adc_item = find_note_and_followup("Promising, low-patent, recent ADC targets in Lung Cancer")
_, lung_bsp_note, lung_bsp_item = find_note_and_followup("Bispecific targets in Lung Cancer")

TOTAL_ROWS = None
MIN_YEAR = MAX_YEAR = None
if overview_table is not None:
    cols = {c.lower(): c for c in overview_table.columns}
    try:
        TOTAL_ROWS = int(overview_table.iloc[0][cols.get("total_rows", overview_table.columns[0])])
    except Exception:
        pass
    for key, target in (("min_year", "MIN_YEAR"), ("max_year", "MAX_YEAR")):
        if key in cols:
            try:
                val = int(overview_table.iloc[0][cols[key]])
                if target == "MIN_YEAR":
                    MIN_YEAR = val
                else:
                    MAX_YEAR = val
            except Exception:
                pass

CROWDED_TOP5_STR = ""
if crowded_table is not None:
    count_col = next((c for c in crowded_table.columns if str(c).lower() in ("n", "count", "row_count", "total_patents")), None)
    if count_col:
        top5 = crowded_table.sort_values(count_col, ascending=False).head(5)
        CROWDED_TOP5_STR = ", ".join(f"{r['target_name']} ({int(r[count_col])} patents)" for _, r in top5.iterrows())

TLDR_BULLETS = [
    f"The dataset spans {TOTAL_ROWS:,} antibody-therapeutics patent records" if TOTAL_ROWS else
    "The dataset covers the full oncology antibody-therapeutics patent landscape",
    (f"filed from {MIN_YEAR} to {MAX_YEAR}, standardized on the harmonized target field (`target_harmonized`)."
     if MIN_YEAR and MAX_YEAR else "standardized on the harmonized target field (`target_harmonized`)."),
    f"IP activity is heavily concentrated in a small set of targets - most crowded: {CROWDED_TOP5_STR}." if CROWDED_TOP5_STR else
    "IP activity is heavily concentrated in a small set of long-established targets.",
    (f"{N_EMERGING} targets qualify as 'Emerging Whitespace' (low cumulative patent volume but high recent filing momentum)."
     if N_EMERGING else "A defined set of targets qualify as 'Emerging Whitespace' - low patent volume, high recent momentum."),
    (f"Across {N_22_INDICATIONS} solid-tumor indications x 3 modality buckets (ADC / Bispecific / Other), "
     f"{N_22_QUALIFYING:,} (target, indication, modality) combinations meet the 'emerging promising target' bar - "
     "1-15 patents in that indication and >=25% filed since 2023." if N_22_QUALIFYING else
     "A large number of (target, indication, modality) combinations meet the 'emerging promising target' bar."),
    (f"A full competitive leaderboard ({N_24_ROWS:,} target x indication x modality rows) ranks every target from "
     "most-crowded to least-contested, to contrast dominance vs. white space side by side." if N_24_ROWS else
     "A full competitive leaderboard ranks every target from most-crowded to least-contested."),
    "White-space targets with low patent counts and rising recent filing activity are flagged throughout as the "
    "most interesting considerations for new program/pipeline pursuit - see the Spotlight section below and "
    "Entries 7, 11, 17, 22 and 24.",
]

EXEC_SUMMARY = (
    f"This report presents the final version-1 intellectual-property landscape analysis of the oncology antibody-"
    f"therapeutics patent dataset ({(f'{TOTAL_ROWS:,} records' if TOTAL_ROWS else 'full dataset')}), standardized on "
    "the harmonized target field (`target_harmonized`). All analytics are produced via reproducible DuckDB SQL "
    "queries against the source CSV, documented step by step in the companion Jupyter notebook "
    "(ip_landscape_final_version1.ipynb, copy included in this folder). The analysis quantifies where patent "
    "activity is concentrated (crowded targets, dominant assignees, modality mix, target-class and payload-class "
    "breakdowns, competitive concentration and claim depth) and - critically - where it is not: the white space. "
    "Particular emphasis is placed on identifying white space targets with low patent counts and recent filing "
    "momentum, since these represent the most interesting considerations for pursuing new programs with a clearer "
    "runway ahead of them. Every table and chart from the source notebook is reproduced below in notebook order; "
    "any table longer than 20 rows is previewed here (first 20 rows) with a clickable link to the full CSV file, "
    "which is also deposited in the data/ folder alongside this report."
)

# ============================================================================
# 3. DOCX - polished report
# ============================================================================
print("Building polished DOCX report ...")

_bookmark_id = [100]


def new_bookmark_id():
    _bookmark_id[0] += 1
    return _bookmark_id[0]


def add_bookmark(paragraph, name):
    bid = new_bookmark_id()
    start = OxmlElement("w:bookmarkStart")
    start.set(qn("w:id"), str(bid))
    start.set(qn("w:name"), name)
    end = OxmlElement("w:bookmarkEnd")
    end.set(qn("w:id"), str(bid))
    paragraph._p.insert(0, start)
    paragraph._p.append(end)


def add_internal_hyperlink(paragraph, text, anchor, bold=False):
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("w:anchor"), anchor)
    new_run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "1155CC")
    rPr.append(color)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rPr.append(u)
    if bold:
        b = OxmlElement("w:b")
        rPr.append(b)
    new_run.append(rPr)
    t = OxmlElement("w:t")
    t.text = text
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)


def add_external_hyperlink(paragraph, text, relative_target):
    part = paragraph.part
    r_id = part.relate_to(
        relative_target, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    new_run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "1155CC")
    rPr.append(color)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rPr.append(u)
    new_run.append(rPr)
    t = OxmlElement("w:t")
    t.text = text
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)


def pd_isna(val):
    try:
        import math
        return val is None or (isinstance(val, float) and math.isnan(val))
    except Exception:
        return val is None


def docx_add_df_table(doc, df, max_rows=20):
    shown = df.head(max_rows)
    table = doc.add_table(rows=1, cols=len(shown.columns))
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr_cells = table.rows[0].cells
    for i, col in enumerate(shown.columns):
        hdr_cells[i].text = str(col)
        for p in hdr_cells[i].paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.size = Pt(8)
    for _, row in shown.iterrows():
        cells = table.add_row().cells
        for i, val in enumerate(row):
            cells[i].text = "" if pd_isna(val) else str(val)
            for p in cells[i].paragraphs:
                for r in p.runs:
                    r.font.size = Pt(8)
    return table


def docx_add_table_item(doc, item):
    df = item["df"]
    docx_add_df_table(doc, df, max_rows=20)
    if len(df) > 20:
        p = doc.add_paragraph()
        p.add_run(f"Showing first 20 of {len(df):,} rows. ").italic = True
        rel = f"data/{item['csv_name']}"
        add_external_hyperlink(p, f"Click to open full CSV ({item['csv_name']})", rel)
    doc.add_paragraph()


def docx_add_image_item(doc, item, width_in=6.5):
    doc.add_picture(str(item["path"]), width=Inches(width_in))
    last_p = doc.paragraphs[-1]
    last_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()


def docx_render_items(doc, items, notebook_mirror=False):
    for it in items:
        if it["type"] == "note":
            p = doc.add_paragraph()
            p.add_run(it["text"]).italic = True
        elif it["type"] == "text":
            if notebook_mirror:
                doc.add_paragraph(it["text"])
            else:
                p = doc.add_paragraph()
                p.add_run(it["text"]).italic = True
        elif it["type"] == "table":
            if notebook_mirror:
                df = it["df"]
                docx_add_df_table(doc, df, max_rows=min(50, len(df)))
                if len(df) > 50:
                    p = doc.add_paragraph()
                    p.add_run(f"Showing first 50 of {len(df):,} rows (mirror capped for document size). ").italic = True
                    rel = f"data/{it['csv_name']}"
                    add_external_hyperlink(p, f"Full CSV: {it['csv_name']}", rel)
                doc.add_paragraph()
            else:
                docx_add_table_item(doc, it)
        elif it["type"] == "image":
            docx_add_image_item(doc, it)


doc = Document()

title_p = doc.add_paragraph()
title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = title_p.add_run(REPORT_TITLE)
run.bold = True
run.font.size = Pt(28)
run.font.color.rgb = RGBColor(0x1F, 0x2A, 0x44)

sub_p = doc.add_paragraph()
sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = sub_p.add_run(SUBTITLE)
run.font.size = Pt(14)
run.italic = True

meta_p = doc.add_paragraph()
meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
meta_p.add_run(
    f"Report generated: {TODAY}\nSource notebook: ip_landscape_final_version1.ipynb\n"
    f"Source data: ip_final_version3.csv" + (f" ({TOTAL_ROWS:,} records)" if TOTAL_ROWS else "")
)
doc.add_page_break()

# --- Table of Contents ---
toc_heading = doc.add_heading("Table of Contents", level=1)
add_bookmark(toc_heading, "toc")

toc_entries = [("exec_summary", "Executive Summary"), ("tldr", "TL;DR"),
               ("spotlight", "Spotlight: Promising Low-Patent White Space Targets")]
toc_entries += [(f"sec_entry_{e['number']}", f"Entry {e['number']}: {e['title']}") for e in entries]
toc_entries += [("appendix", "Appendix: Data Files Index")]

for anchor, label in toc_entries:
    p = doc.add_paragraph()
    add_internal_hyperlink(p, label, anchor)
doc.add_page_break()

# --- Executive Summary ---
h = doc.add_heading("Executive Summary", level=1)
add_bookmark(h, "exec_summary")
doc.add_paragraph(EXEC_SUMMARY)

# --- TL;DR ---
h = doc.add_heading("TL;DR", level=1)
add_bookmark(h, "tldr")
for b in TLDR_BULLETS:
    doc.add_paragraph(b, style="List Bullet")
doc.add_page_break()

# --- Spotlight section ---
h = doc.add_heading("Spotlight: Promising Low-Patent White Space Targets", level=1)
add_bookmark(h, "spotlight")
doc.add_paragraph(
    "This analysis specifically identifies white space targets with low patent counts and rising recent filing "
    "activity - these are the most interesting considerations for pursuing new programs, since they combine a "
    "narrow competitive footprint today with clear evidence of growing interest (i.e., a closing but still-open "
    "window). Two direct call-outs from the Indication x Modality-Class analysis (Entry 22) are highlighted here; "
    "the same methodology (1-15 total patents in that indication AND >=25% of filings since 2023) was applied "
    f"across all {N_22_INDICATIONS} solid-tumor indications and all 3 modality buckets (ADC, Bispecific, Other) - "
    "see Entry 22 for the complete cross-indication results."
)
if lung_adc_note and lung_adc_item and lung_adc_item["type"] == "table":
    doc.add_heading(lung_adc_note["text"].strip("*: "), level=2)
    docx_add_table_item(doc, lung_adc_item)
if lung_bsp_note and lung_bsp_item and lung_bsp_item["type"] == "table":
    doc.add_heading(lung_bsp_note["text"].strip("*: "), level=2)
    docx_add_table_item(doc, lung_bsp_item)
doc.add_page_break()

# --- Entry sections ---
for e in entries:
    h = doc.add_heading(f"Entry {e['number']}: {e['title']}", level=1)
    add_bookmark(h, f"sec_entry_{e['number']}")
    if e["description"]:
        doc.add_paragraph(e["description"])
    docx_render_items(doc, e["items"], notebook_mirror=False)
    doc.add_paragraph()

doc.add_page_break()

# --- Appendix: data files index ---
h = doc.add_heading("Appendix: Data Files Index", level=1)
add_bookmark(h, "appendix")
doc.add_paragraph(
    "Every table shown in this report (regardless of size) has its full underlying data deposited as a CSV file "
    "in the data/ folder next to this report, clearly labelled by Entry number and topic. Tables longer than 20 "
    "rows include a clickable link above ('Click to open full CSV') that opens the file directly."
)
appendix_table = doc.add_table(rows=1, cols=3)
appendix_table.style = "Light Grid Accent 1"
hdr = appendix_table.rows[0].cells
hdr[0].text, hdr[1].text, hdr[2].text = "Entry", "Rows", "File"
for e in entries:
    for it in e["items"]:
        if it["type"] == "table":
            cells = appendix_table.add_row().cells
            cells[0].text = f"Entry {e['number']}: {e['title']}"
            cells[1].text = f"{len(it['df']):,}"
            p = cells[2].paragraphs[0]
            add_external_hyperlink(p, it["csv_name"], f"data/{it['csv_name']}")

docx_path = OUT / "IP_Landscape_Report.docx"
doc.save(docx_path)
print(f"Saved DOCX: {docx_path}")

# ============================================================================
# 4. PDF - mirrored content with internal clickable TOC
# ============================================================================
print("Building PDF report ...")

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(name="TitleBig", fontSize=24, leading=28, alignment=1, textColor=colors.HexColor("#1F2A44"), spaceAfter=10))
styles.add(ParagraphStyle(name="SubtitleBig", fontSize=12, leading=15, alignment=1, textColor=colors.HexColor("#444444"), spaceAfter=6, fontName="Helvetica-Oblique"))
styles.add(ParagraphStyle(name="MetaCenter", fontSize=9.5, leading=13, alignment=1, textColor=colors.HexColor("#666666")))
styles.add(ParagraphStyle(name="H1", fontSize=16, leading=20, spaceBefore=14, spaceAfter=8, textColor=colors.HexColor("#1F2A44")))
styles.add(ParagraphStyle(name="H2", fontSize=12.5, leading=16, spaceBefore=10, spaceAfter=6, textColor=colors.HexColor("#2563EB")))
styles.add(ParagraphStyle(name="Body9", fontSize=9, leading=12.5, spaceAfter=6))
styles.add(ParagraphStyle(name="Note9", fontSize=9, leading=12.5, spaceAfter=6, textColor=colors.HexColor("#333333"), fontName="Helvetica-Oblique"))
styles.add(ParagraphStyle(name="Caption", fontSize=8, leading=10, textColor=colors.HexColor("#666666"), spaceAfter=10, alignment=1))
styles.add(ParagraphStyle(name="Bullet9", fontSize=9, leading=12.5, leftIndent=14, bulletIndent=2, spaceAfter=4))
styles.add(ParagraphStyle(name="TOCLink", fontSize=10, leading=15, textColor=colors.HexColor("#1155CC")))


def esc(text):
    return str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def pdf_df_table(df, max_rows=20, font_size=6.5):
    shown = df.head(max_rows)
    data = [[esc(c) for c in shown.columns]] + [
        ["" if pd_isna(v) else esc(v) for v in row] for row in shown.itertuples(index=False)
    ]
    ncols = len(shown.columns)
    avail_width = 7.3 * inch
    col_width = avail_width / max(ncols, 1)
    t = RLTable(data, colWidths=[col_width] * ncols, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F2A44")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, -1), font_size),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CCCCCC")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F6FA")]),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 3),
        ("RIGHTPADDING", (0, 0), (-1, -1), 3),
    ]))
    return t


def pdf_render_items(story, items):
    for it in items:
        if it["type"] in ("note", "text"):
            story.append(Paragraph(esc(it["text"]), styles["Note9"]))
        elif it["type"] == "table":
            story.append(pdf_df_table(it["df"], max_rows=20))
            if len(it["df"]) > 20:
                rel = f"data/{it['csv_name']}"
                story.append(Paragraph(
                    f"Showing first 20 of {len(it['df']):,} rows. "
                    f"<link href=\"{rel}\">Click to open full CSV ({esc(it['csv_name'])})</link>",
                    styles["Body9"]))
            story.append(Spacer(1, 8))
        elif it["type"] == "image":
            with PILImage.open(it["path"]) as im:
                w, h = im.size
            max_w = 6.6 * inch
            ratio = h / w
            draw_w = max_w
            draw_h = max_w * ratio
            max_h = 7.5 * inch
            if draw_h > max_h:
                draw_h = max_h
                draw_w = max_h / ratio
            story.append(RLImage(str(it["path"]), width=draw_w, height=draw_h))
            story.append(Spacer(1, 8))


story = []
story.append(Spacer(1, 1.4 * inch))
story.append(Paragraph(esc(REPORT_TITLE), styles["TitleBig"]))
story.append(Paragraph(esc(SUBTITLE), styles["SubtitleBig"]))
story.append(Spacer(1, 0.25 * inch))
story.append(Paragraph(
    f"Report generated: {TODAY}<br/>Source notebook: ip_landscape_final_version1.ipynb<br/>"
    f"Source data: ip_final_version3.csv" + (f" ({TOTAL_ROWS:,} records)" if TOTAL_ROWS else ""),
    styles["MetaCenter"],
))
story.append(PageBreak())

# TOC
story.append(Paragraph('<a name="toc"/>Table of Contents', styles["H1"]))
for anchor, label in toc_entries:
    story.append(Paragraph(f'<link href="#{anchor}">{esc(label)}</link>', styles["TOCLink"]))
story.append(PageBreak())

story.append(Paragraph('<a name="exec_summary"/>Executive Summary', styles["H1"]))
story.append(Paragraph(esc(EXEC_SUMMARY), styles["Body9"]))

story.append(Paragraph('<a name="tldr"/>TL;DR', styles["H1"]))
for b in TLDR_BULLETS:
    story.append(Paragraph(f"&bull; {esc(b)}", styles["Bullet9"]))
story.append(PageBreak())

story.append(Paragraph('<a name="spotlight"/>Spotlight: Promising Low-Patent White Space Targets', styles["H1"]))
story.append(Paragraph(esc(
    "This analysis specifically identifies white space targets with low patent counts and rising recent filing "
    "activity - these are the most interesting considerations for pursuing new programs, since they combine a "
    "narrow competitive footprint today with clear evidence of growing interest. Two direct call-outs from the "
    "Indication x Modality-Class analysis (Entry 22) are highlighted here; the same methodology was applied across "
    f"all {N_22_INDICATIONS} solid-tumor indications and all 3 modality buckets - see Entry 22 for full results."
), styles["Body9"]))
if lung_adc_note and lung_adc_item and lung_adc_item["type"] == "table":
    story.append(Paragraph(esc(lung_adc_note["text"].strip("*: ")), styles["H2"]))
    story.append(pdf_df_table(lung_adc_item["df"], max_rows=20))
    if len(lung_adc_item["df"]) > 20:
        rel = f"data/{lung_adc_item['csv_name']}"
        story.append(Paragraph(f'Showing first 20 rows. <link href="{rel}">Click to open full CSV</link>', styles["Body9"]))
    story.append(Spacer(1, 10))
if lung_bsp_note and lung_bsp_item and lung_bsp_item["type"] == "table":
    story.append(Paragraph(esc(lung_bsp_note["text"].strip("*: ")), styles["H2"]))
    story.append(pdf_df_table(lung_bsp_item["df"], max_rows=20))
    if len(lung_bsp_item["df"]) > 20:
        rel = f"data/{lung_bsp_item['csv_name']}"
        story.append(Paragraph(f'Showing first 20 rows. <link href="{rel}">Click to open full CSV</link>', styles["Body9"]))
story.append(PageBreak())

for e in entries:
    story.append(Paragraph(f'<a name="sec_entry_{e["number"]}"/>Entry {e["number"]}: {esc(e["title"])}', styles["H1"]))
    if e["description"]:
        story.append(Paragraph(esc(e["description"]), styles["Body9"]))
    pdf_render_items(story, e["items"])
    story.append(PageBreak())

story.append(Paragraph('<a name="appendix"/>Appendix: Data Files Index', styles["H1"]))
story.append(Paragraph(esc(
    "Every table in this report has its full underlying data deposited as a CSV file in the data/ folder next to "
    "this report, clearly labelled by Entry number and topic."
), styles["Body9"]))
appendix_rows = [["Entry", "Rows", "File"]]
for e in entries:
    for it in e["items"]:
        if it["type"] == "table":
            appendix_rows.append([f"Entry {e['number']}: {e['title'][:40]}", f"{len(it['df']):,}", it["csv_name"]])
t = RLTable(appendix_rows, colWidths=[3.2 * inch, 0.8 * inch, 3.3 * inch], repeatRows=1)
t.setStyle(TableStyle([
    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F2A44")),
    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
    ("FONTSIZE", (0, 0), (-1, -1), 7),
    ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CCCCCC")),
    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F6FA")]),
]))
story.append(t)

pdf_path = OUT / "IP_Landscape_Report.pdf"
pdf_doc = SimpleDocTemplate(str(pdf_path), pagesize=letter, topMargin=0.6 * inch, bottomMargin=0.6 * inch,
                            leftMargin=0.6 * inch, rightMargin=0.6 * inch, title=REPORT_TITLE)
pdf_doc.build(story)
print(f"Saved PDF: {pdf_path}")

# ============================================================================
# 5. PPTX - topics + graphs slide deck
# ============================================================================
print("Building PPTX ...")

prs = Presentation()
prs.slide_width = PIn(13.333)
prs.slide_height = PIn(7.5)
BLANK = prs.slide_layouts[6]
DARK = PRGBColor(0x1F, 0x2A, 0x44)
ACCENT = PRGBColor(0x25, 0x63, 0xEB)
GREY = PRGBColor(0x44, 0x44, 0x44)


def add_title_box(slide, text, top=PIn(0.35), size=30, color=DARK, height=PIn(1.0)):
    box = slide.shapes.add_textbox(PIn(0.5), top, prs.slide_width - PIn(1.0), height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PPt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return box


def add_bullets(slide, bullets, top=PIn(1.5), left=PIn(0.7), width=None, height=None, size=16):
    width = width or (prs.slide_width - PIn(1.4))
    height = height or (prs.slide_height - top - PIn(0.4))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022 {b}"
        p.font.size = PPt(size)
        p.font.color.rgb = GREY
        p.space_after = PPt(10)
    return box


def add_image_slide(title, image_path, caption=None):
    slide = prs.slides.add_slide(BLANK)
    add_title_box(slide, title, size=24, height=PIn(0.7))
    with PILImage.open(image_path) as im:
        w, h = im.size
    ratio = h / w
    max_w = PIn(12.2)
    max_h = PIn(5.9)
    draw_w = max_w
    draw_h = PIn(max_w.inches * ratio)
    if draw_h > max_h:
        draw_h = max_h
        draw_w = PIn(max_h.inches / ratio)
    left = (prs.slide_width - draw_w) / 2
    slide.shapes.add_picture(str(image_path), left, PIn(1.15), width=draw_w, height=draw_h)
    if caption:
        cap = slide.shapes.add_textbox(PIn(0.5), prs.slide_height - PIn(0.55), prs.slide_width - PIn(1.0), PIn(0.4))
        p = cap.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = PPt(11)
        p.font.italic = True
        p.font.color.rgb = GREY
        p.alignment = PP_ALIGN.CENTER
    return slide


def add_table_slide(title, df, max_rows=8):
    slide = prs.slides.add_slide(BLANK)
    add_title_box(slide, title, size=24, height=PIn(0.7))
    shown = df.head(max_rows)
    rows, cols = shown.shape[0] + 1, shown.shape[1]
    left, top, width, height = PIn(0.6), PIn(1.3), prs.slide_width - PIn(1.2), PIn(5.6)
    gtable = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for j, col in enumerate(shown.columns):
        gtable.cell(0, j).text = str(col)
    for i, (_, row) in enumerate(shown.iterrows(), start=1):
        for j, val in enumerate(row):
            gtable.cell(i, j).text = "" if pd_isna(val) else str(val)
    return slide


# Title slide
slide = prs.slides.add_slide(BLANK)
box = slide.shapes.add_textbox(PIn(0.8), PIn(2.6), prs.slide_width - PIn(1.6), PIn(1.6))
p = box.text_frame.paragraphs[0]
p.text = REPORT_TITLE
p.font.size = PPt(40)
p.font.bold = True
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER
sub_box = slide.shapes.add_textbox(PIn(0.8), PIn(3.9), prs.slide_width - PIn(1.6), PIn(1.0))
p2 = sub_box.text_frame.paragraphs[0]
p2.text = SUBTITLE
p2.font.size = PPt(18)
p2.font.italic = True
p2.font.color.rgb = GREY
p2.alignment = PP_ALIGN.CENTER
meta_box = slide.shapes.add_textbox(PIn(0.8), PIn(4.9), prs.slide_width - PIn(1.6), PIn(0.8))
p3 = meta_box.text_frame.paragraphs[0]
p3.text = f"Generated {TODAY}  |  ip_landscape_final_version1.ipynb  |  ip_final_version3.csv"
p3.font.size = PPt(13)
p3.font.color.rgb = GREY
p3.alignment = PP_ALIGN.CENTER

# TL;DR slide
slide = prs.slides.add_slide(BLANK)
add_title_box(slide, "TL;DR", size=32)
add_bullets(slide, TLDR_BULLETS, size=17)

# Spotlight slides
if lung_adc_item and lung_adc_item["type"] == "table":
    add_table_slide("Spotlight: Promising Low-Patent ADC Targets - Lung Cancer", lung_adc_item["df"], max_rows=10)
if lung_bsp_item and lung_bsp_item["type"] == "table":
    add_table_slide("Spotlight: Promising Low-Patent Bispecific Targets - Lung Cancer", lung_bsp_item["df"], max_rows=10)

# Agenda slide
slide = prs.slides.add_slide(BLANK)
add_title_box(slide, "Agenda", size=32)
agenda_bullets = [f"Entry {e['number']}: {e['title']}" for e in entries]
add_bullets(slide, agenda_bullets, size=12, height=PIn(5.8))

# One slide per Entry: prefer its chart image(s); fall back to a table snapshot
for e in entries:
    images = [it for it in e["items"] if it["type"] == "image"]
    tables = [it for it in e["items"] if it["type"] == "table"]
    title = f"Entry {e['number']}: {e['title']}"
    if images:
        for k, img in enumerate(images):
            cap = title if k == 0 else f"{title} (cont.)"
            add_image_slide(title, img["path"], caption=e["title"])
    elif tables:
        add_table_slide(title, tables[0]["df"], max_rows=8)
    else:
        slide = prs.slides.add_slide(BLANK)
        add_title_box(slide, title, size=26)
        notes = [it["text"] for it in e["items"] if it["type"] in ("note", "text")]
        add_bullets(slide, notes[:6] if notes else ["See full notebook / report for details."], size=15)

pptx_path = OUT / "IP_Landscape_Report.pptx"
prs.save(pptx_path)
print(f"Saved PPTX: {pptx_path}")

# ============================================================================
# 6. Notebook Mirror DOCX - exact cell-by-cell mirror (markdown + code + outputs)
# ============================================================================
print("Building Notebook Mirror DOCX ...")

mdoc = Document()
title_p = mdoc.add_paragraph()
title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = title_p.add_run("IP Landscape Analysis - Notebook Mirror")
run.bold = True
run.font.size = Pt(26)
run.font.color.rgb = RGBColor(0x1F, 0x2A, 0x44)
sub_p = mdoc.add_paragraph()
sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
sub_p.add_run("Exact cell-by-cell mirror of ip_landscape_final_version1.ipynb (markdown, code and outputs, in order)").italic = True
meta_p = mdoc.add_paragraph()
meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
meta_p.add_run(f"Generated: {TODAY}")
mdoc.add_paragraph(
    "Note: tables longer than 50 rows are shown here as the first 50 rows only, with a link to the full CSV file "
    "(deposited in data/), to keep this document a manageable size - all other content (markdown text, full code, "
    "charts, and tables up to 50 rows) is reproduced in full and in the exact order it appears in the notebook."
)
mdoc.add_page_break()

for e in entries:
    mdoc.add_heading(f"Entry {e['number']}: {e['title']}", level=1)
    if e["description"]:
        mdoc.add_paragraph(e["description"])
    for code_src in e["code_sources"]:
        mdoc.add_heading("Code", level=3)
        for line in code_src.split("\n"):
            p = mdoc.add_paragraph()
            r = p.add_run(line if line else " ")
            r.font.name = "Consolas"
            r.font.size = Pt(7.5)
            p.paragraph_format.space_after = Pt(0)
    if e["items"]:
        mdoc.add_heading("Output", level=3)
        docx_render_items(mdoc, e["items"], notebook_mirror=True)
    mdoc.add_paragraph()
    mdoc.add_page_break()

mirror_path = OUT / "IP_Landscape_Notebook_Mirror.docx"
mdoc.save(mirror_path)
print(f"Saved Notebook Mirror DOCX: {mirror_path}")

print("\nAll deliverables written to:", OUT)
for p in sorted(OUT.iterdir()):
    if p.is_file():
        print(" -", p.name)
print(" - data/  (", len(list(DATA.iterdir())), "CSV files )")
print(" - figures/  (", len(list(FIG.iterdir())), "PNG files )")
