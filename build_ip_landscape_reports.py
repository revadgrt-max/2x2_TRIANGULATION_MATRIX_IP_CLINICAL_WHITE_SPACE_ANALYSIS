"""
Build the DOCX, PDF, and PPTX deliverables for the IP Antibody Therapeutics
Landscape Analysis, reusing every table/figure produced by
build_ip_landscape_report.py (imported as a module so all analytics stay
perfectly consistent with the source notebook).
"""

from datetime import date

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Inches, Pt, RGBColor

from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table as RLTable,
    TableStyle, PageBreak, HRFlowable,
)

print("Importing analytics module (recomputes data + figures) ...")
import build_ip_landscape_report as ipl  # noqa: E402  (re-executes the full pipeline)

OUT = ipl.OUT
DATA = ipl.DATA
FIG = ipl.FIG
REPORT_TITLE = "IP Antibody Therapeutics Landscape Analysis"
SUBTITLE = "Patent Landscape, Competitive Crowding & White Space Analysis - Oncology Antibody Therapeutics"
TODAY = date(2026, 7, 25).strftime("%B %d, %Y")

TOTAL_ROWS = int(ipl.total_rows_df["total_rows"].iloc[0])
CROWDED_TOP10 = ipl.top_targets_df.dropna(subset=["target_name"]).head(10)
EMERGING_TOP10 = ipl.emerging_deepdive_df.sort_values("total_patents", ascending=False).head(10)
WHITESPACE_INDICATIONS = ipl.whitespace_per_indication_df
N_EMERGING = len(ipl.emerging_deepdive_df)
N_ZERO_CELLS_G = int((ipl.target_modality_indication_grid_df["patent_count"] == 0).sum())
N_TOTAL_CELLS_G = len(ipl.target_modality_indication_grid_df)
N_ZERO_CELLS_H = int((ipl.solid_indication_target_grid_df["patent_count"] == 0).sum())
N_TOTAL_CELLS_H = len(ipl.solid_indication_target_grid_df)

CROWDED_LIST_STR = ", ".join(f"{r.target_name} ({int(r.n)} patents)" for r in CROWDED_TOP10.itertuples())
STANDOUT_WHITESPACE = ", ".join(sorted(set(
    t.strip() for targets in WHITESPACE_INDICATIONS["whitespace_targets"]
    for t in str(targets).split(",") if "none" not in str(targets).lower()
)))

TLDR_BULLETS = [
    f"The dataset covers {TOTAL_ROWS:,} antibody-therapeutics patent records in oncology, spanning filings from "
    f"{int(ipl.overview_df['min_year'].iloc[0])} to {int(ipl.overview_df['max_year'].iloc[0])}.",
    f"IP is heavily concentrated in a small set of targets: PD-1, HER2, PD-L1, EGFR and CD3 alone account for "
    f"{int(CROWDED_TOP10['n'].head(5).sum()):,} patents ({100*CROWDED_TOP10['n'].head(5).sum()/TOTAL_ROWS:.1f}% of all filings).",
    f"{N_EMERGING} targets qualify as 'Emerging Whitespace' (under 40 total patents but >=30% of activity filed since 2023) - "
    "led by CCR8, DLL3, ROR1, BCMA x CD3 and GPC3.",
    f"Across the 20 most-established targets x 12 leading solid-tumor indications ({N_TOTAL_CELLS_H} combinations), "
    f"{N_ZERO_CELLS_H} target x indication pairs ({100*N_ZERO_CELLS_H/N_TOTAL_CELLS_H:.0f}%) have zero filed patents - "
    "genuine white space.",
    f"Two targets stand out as unclaimed across the widest range of solid tumor indications: {STANDOUT_WHITESPACE}.",
    "The analysis surfaces clear white space targets - both at the individual-target level (Emerging Whitespace "
    "quadrant) and at the target x indication level (unclaimed combinations) - giving a data-driven shortlist for "
    "pipeline prioritization and freedom-to-operate scouting.",
]

EXEC_SUMMARY = (
    f"This report presents a comprehensive intellectual-property landscape analysis of {TOTAL_ROWS:,} oncology "
    "antibody-therapeutics patent records, built entirely from the FINAL_master_plus_mab_addendum.csv dataset using "
    "reproducible DuckDB SQL queries (documented in the companion Jupyter notebook, ip_lanscape.ipynb). The analysis "
    "quantifies where patent activity is concentrated (crowded targets, dominant assignees, modality mix), how that "
    "activity is trending over time, and - critically - where it is not: the white space. "
    f"The most crowded targets in the dataset are {CROWDED_LIST_STR}, reflecting mature, heavily-litigated antibody "
    "franchises (checkpoint inhibitors, HER2- and EGFR-directed therapies, and CD3-engaging bispecifics). Layered "
    "against this crowding picture, the analysis has surfaced clear white space targets - both emerging, "
    "low-volume/high-momentum targets (e.g., CCR8, DLL3, ROR1) and specific target x indication combinations "
    "(most notably Mesothelin and OX40 across multiple solid tumor indications) that remain largely unclaimed "
    "relative to the 20 most-established targets in the space. Every table and chart from the source notebook is "
    "reproduced below; any table exceeding 10 rows is summarized here and provided in full as a linked CSV file "
    "(and consolidated into a single multi-sheet Excel workbook) in the accompanying data/ folder."
)

RECOMMENDATIONS = [
    "Prioritize freedom-to-operate (FTO) review before pursuing PD-1, HER2, PD-L1, EGFR or CD3 as primary novelty "
    "claims - these are saturated with 60+ unique assignees each in many cases.",
    "Fast-track internal programs against CCR8, DLL3, ROR1, GPC3 and Nectin-4 - all show low cumulative patent "
    "volume but >45% of filings since 2023, indicating a narrow and closing window of opportunity.",
    "Evaluate Mesothelin- and OX40-directed candidates specifically in breast, colorectal, gastric/GEJ, "
    "glioblastoma, liver/HCC and prostate indications, where the target x indication grid shows zero competing "
    "filings among the top 20 targets.",
    "Treat competitive concentration as a second filter: several emerging targets are already dominated by a single "
    "assignee holding a large majority of filings - validate freedom-to-operate against that specific assignee, "
    "not just the target as a whole.",
    "Revisit this analysis on a recurring (e.g., quarterly) cadence, since the whitespace/crowded classification is "
    "explicitly time-sensitive (based on % of filings in the last 3 years).",
]

STANDOUT_SUGGESTIONS = [
    "Interactive dashboard: publish the target x indication and target x modality grids as an interactive "
    "Plotly/Dash or Tableau view so stakeholders can filter by therapeutic area or assignee themselves.",
    "One-page executive infographic: a single visual page (crowded-vs-whitespace quadrant + top 5 white space "
    "callouts) for leadership-level consumption, distinct from the full technical report.",
    "Competitive positioning matrix: an assignee x target bubble chart (bubble size = patent count, color = filing "
    "recency) to show not just which targets are crowded, but who is crowding them.",
    "Confidence/provenance labeling: since indication and target resolution rely on parsed/derived fields, add a "
    "visible confidence tier (High/Medium/Low) per finding, tied to the existing resolution_confidence field.",
    "FTO risk scoring: combine assignee concentration, claim scope (independent claim count) and recency into a "
    "single composite 'competitive risk score' per target to make prioritization even more actionable.",
    "Living document versioning: timestamp and changelog each re-run of this pipeline so trend shifts (e.g., a "
    "target moving from 'Emerging Whitespace' to 'Hot/Crowded') are tracked over time, not just snapshotted.",
    "Direct data drill-down links: since every large table already has a linked CSV/Excel export, consider adding "
    "QR codes on the printed/PDF version that jump straight to the relevant data file or dashboard.",
]

DATA_QUALITY_NOTES = [
    "Indication and target fields are partially derived/parsed from free text; records where the parser could not "
    "resolve a specific tumor type are grouped under the generic label 'oncology TA query hit' and excluded from "
    "indication-level whitespace analysis.",
    "Target-level analysis uses target_corrected where available, falling back to the raw target field; ~3,000 "
    "records had no resolvable target and are excluded from target-level tables.",
    "Assignee names are not fully normalized (e.g., corporate entity variants and non-English co-listings appear as "
    "distinct strings), which can slightly understate true assignee concentration.",
    "'Recent' filing activity is defined as filing_year >= 2023; whitespace/crowded quadrant classifications are "
    "therefore sensitive to the current date and should be refreshed periodically.",
    "The solid-tumor indication grid explicitly excludes heme/liquid tumor keywords (leukemia, lymphoma, myeloma, "
    "CLL, AML, ALL, MDS, Hodgkin, Waldenstrom, 'blood') to isolate solid-tumor-specific white space.",
]

# ============================================================================
# Shared helpers to decide preview size + link text per table
# ============================================================================

def preview_rows(df, n=10):
    return df.head(n)


def linked_tables():
    """Return only the entries that were exported to CSV (row count > 10)."""
    return [t for t in ipl.TABLE_REGISTRY if t["csv_path"] is not None]


# ============================================================================
# DOCX BUILD
# ============================================================================
print("Building DOCX ...")


def add_hyperlink(paragraph, text, target_path):
    """Add a relative-path hyperlink (to a local file) into a docx paragraph."""
    part = paragraph.part
    r_id = part.relate_to(
        target_path, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    new_run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "1155CC")
    rPr.append(color)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rPr.append(u)
    new_run.append(rPr)
    t = OxmlElement("w:t")
    t.text = text
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)
    return hyperlink


def docx_add_df_table(doc, df, max_rows=10):
    shown = df.head(max_rows)
    table = doc.add_table(rows=1, cols=len(shown.columns))
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr_cells = table.rows[0].cells
    for i, col in enumerate(shown.columns):
        hdr_cells[i].text = str(col)
        for p in hdr_cells[i].paragraphs:
            for r in p.runs:
                r.bold = True
    for _, row in shown.iterrows():
        cells = table.add_row().cells
        for i, val in enumerate(row):
            cells[i].text = "" if pd_isna(val) else str(val)
    return table


def pd_isna(val):
    try:
        import math
        return val is None or (isinstance(val, float) and math.isnan(val))
    except Exception:
        return val is None


def docx_add_table_section(doc, entry, note=""):
    doc.add_heading(entry["title"], level=2)
    if note:
        doc.add_paragraph(note)
    docx_add_df_table(doc, entry["df"])
    if entry["csv_path"] is not None:
        p = doc.add_paragraph()
        p.add_run(f"Full table ({len(entry['df']):,} rows): ").italic = True
        rel = f"data/{entry['csv_path'].name}"
        add_hyperlink(p, rel, rel)
    doc.add_paragraph()


doc = Document()

# Title page
title_p = doc.add_paragraph()
title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = title_p.add_run(REPORT_TITLE)
run.bold = True
run.font.size = Pt(28)
run.font.color.rgb = RGBColor(0x1F, 0x2A, 0x44)

sub_p = doc.add_paragraph()
sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = sub_p.add_run(SUBTITLE)
run.font.size = Pt(14)
run.italic = True

meta_p = doc.add_paragraph()
meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
meta_p.add_run(f"Report generated: {TODAY}\nSource notebook: ip_lanscape.ipynb\nSource data: FINAL_master_plus_mab_addendum.csv ({TOTAL_ROWS:,} records)")
doc.add_page_break()

doc.add_heading("Executive Summary", level=1)
doc.add_paragraph(EXEC_SUMMARY)

doc.add_heading("TL;DR", level=1)
for b in TLDR_BULLETS:
    doc.add_paragraph(b, style="List Bullet")

doc.add_heading("Most Crowded Targets", level=2)
p = doc.add_paragraph()
p.add_run("Most crowded targets: ").bold = True
p.add_run(CROWDED_LIST_STR)

doc.add_heading("White Space Targets Surfaced", level=2)
p = doc.add_paragraph()
p.add_run(
    "This analysis has surfaced clear white space targets - see the Emerging Whitespace Deep-Dive and White Space "
    f"Targets per Solid Tumor Indication sections below. Standout unclaimed targets: {STANDOUT_WHITESPACE}."
)
doc.add_page_break()

# Section: Methodology
doc.add_heading("1. Dataset Overview & Methodology", level=1)
doc.add_paragraph(
    f"The underlying dataset contains {TOTAL_ROWS:,} patent records related to antibody therapeutics in oncology, "
    f"covering filing years {int(ipl.overview_df['min_year'].iloc[0])}-{int(ipl.overview_df['max_year'].iloc[0])}. "
    "All analytics in this report are produced via DuckDB SQL queries executed directly against the source CSV "
    "(no intermediate database), mirroring the queries documented in the Code Log of ip_lanscape.ipynb."
)
doc.add_picture(str(ipl.fig01_path), width=Inches(6.5))
doc.add_paragraph("Figure 1. Overview dashboard: filing trend, modality mix, top targets, top assignees, indication whitespace and target x modality crowding matrix.", style="Caption")

doc.add_heading("2. Modality Distribution", level=2)
docx_add_df_table(doc, ipl.modality_df, max_rows=10)
doc.add_paragraph()

doc.add_heading("3. Filing Trend by Year", level=2)
docx_add_table_section(doc, next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table04_filing_trend_by_year"))

# Crowded targets / assignees / indications / claim_focus / crowding relevance
for key in ["table01_top25_targets", "table03_top25_assignees", "table02_top25_indication_strings",
            "table05_claim_focus_top15", "table06_crowding_relevance_top15",
            "table07_parsed_indication_whitespace", "table08_target_modality_crowding_matrix_top15"]:
    entry = next(t for t in ipl.TABLE_REGISTRY if t["key"] == key)
    docx_add_table_section(doc, entry)

doc.add_heading("4. Target Whitespace Scoring & Multi-Angle View", level=1)
docx_add_table_section(doc, next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table09_target_whitespace_scoring_top30"))
doc.add_picture(str(ipl.fig02_path), width=Inches(6.5))
doc.add_paragraph("Figure 2. Whitespace quadrant scatter, target x indication crowding heatmap, modality share and indication-resolution rate.", style="Caption")
for key in ["table10A_target_whitespace_summary_all", "table10B_target_modality_share_all",
            "table10C_target_indication_crowding_matrix", "table10D_whitespace_quadrant_classification"]:
    entry = next(t for t in ipl.TABLE_REGISTRY if t["key"] == key)
    docx_add_table_section(doc, entry)

doc.add_heading("5. Emerging Whitespace Deep-Dive", level=1)
doc.add_paragraph(
    f"{N_EMERGING} targets meet the Emerging Whitespace definition (total patents < 40 and >=30% filed since 2023). "
    "The table below shows the top 10 by volume; the full set of "
    f"{N_EMERGING} targets is available in the linked CSV."
)
docx_add_df_table(doc, EMERGING_TOP10[["target_name", "total_patents", "pct_recent", "top_assignee", "dominant_modality"]], max_rows=10)
entry = next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table11E_emerging_whitespace_deepdive")
p = doc.add_paragraph()
p.add_run(f"Full table ({len(entry['df']):,} rows): ").italic = True
rel = f"data/{entry['csv_path'].name}"
add_hyperlink(p, rel, rel)
doc.add_picture(str(ipl.fig03_path), width=Inches(6.5))
doc.add_paragraph("Figure 3. Emerging whitespace volume/momentum, filing trend, assignee concentration and dominant modality share.", style="Caption")
entry = next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table11F_emerging_whitespace_filing_trend")
docx_add_table_section(doc, entry, note="Year-by-year filing counts for every emerging whitespace target (2015-2025).")

doc.add_heading("6. Target x Modality x Indication Whitespace Grid", level=1)
doc.add_paragraph(
    f"Across the top 10 targets, top 5 modalities and top 8 indications ({N_TOTAL_CELLS_G} combinations), "
    f"{N_ZERO_CELLS_G} combinations ({100*N_ZERO_CELLS_G/N_TOTAL_CELLS_G:.0f}%) have zero filed patents."
)
doc.add_picture(str(ipl.fig04a_path), width=Inches(6.5))
doc.add_paragraph("Figure 4a. Small-multiple heatmaps of modality x indication patent counts per target (green/0 = white space, red = crowded).", style="Caption")
doc.add_picture(str(ipl.fig04b_path), width=Inches(6.0))
doc.add_paragraph("Figure 4b. Whitespace breadth (unclaimed modality x indication cells) per target.", style="Caption")
entry = next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table12G_target_modality_indication_grid")
docx_add_table_section(doc, entry)

doc.add_heading("7. White Space Targets Available per Solid Tumor Indication", level=1)
doc.add_paragraph(
    f"Considering the 20 most-established targets against the 12 leading solid tumor indications "
    f"({N_TOTAL_CELLS_H} combinations), {N_ZERO_CELLS_H} combinations ({100*N_ZERO_CELLS_H/N_TOTAL_CELLS_H:.0f}%) "
    f"are entirely unclaimed. Standout white space targets: {STANDOUT_WHITESPACE}."
)
docx_add_df_table(doc, WHITESPACE_INDICATIONS, max_rows=12)
doc.add_picture(str(ipl.fig05_path), width=Inches(6.5))
doc.add_paragraph("Figure 5. Solid tumor indication x target whitespace heatmap and whitespace breadth per indication.", style="Caption")
entry = next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table13H_solid_tumor_indication_target_grid")
p = doc.add_paragraph()
p.add_run(f"Full grid ({len(entry['df']):,} rows): ").italic = True
rel = f"data/{entry['csv_path'].name}"
add_hyperlink(p, rel, rel)
doc.add_paragraph()

doc.add_heading("8. Key Findings & Recommendations", level=1)
for r in RECOMMENDATIONS:
    doc.add_paragraph(r, style="List Bullet")

doc.add_heading("9. Data Quality, Caveats & Methodology Notes", level=1)
for n in DATA_QUALITY_NOTES:
    doc.add_paragraph(n, style="List Bullet")

doc.add_heading("10. Making This Report Stand Out - Suggestions", level=1)
for s in STANDOUT_SUGGESTIONS:
    doc.add_paragraph(s, style="List Bullet")

doc.add_heading("Appendix: Linked Data Files", level=1)
doc.add_paragraph(
    "Every table above with more than 10 rows is exported in full to the data/ folder as an individual CSV, and "
    "all tables (regardless of size) are also consolidated into a single multi-sheet Excel workbook."
)
p = doc.add_paragraph()
p.add_run("Combined workbook: ").italic = True
add_hyperlink(p, f"data/{ipl.xlsx_path.name}", f"data/{ipl.xlsx_path.name}")
appendix_table = doc.add_table(rows=1, cols=3)
appendix_table.style = "Light Grid Accent 1"
hdr = appendix_table.rows[0].cells
hdr[0].text, hdr[1].text, hdr[2].text = "Table", "Rows", "File"
for t in linked_tables():
    cells = appendix_table.add_row().cells
    cells[0].text = t["title"]
    cells[1].text = str(len(t["df"]))
    cells[2].text = f"data/{t['csv_path'].name}"

docx_path = OUT / "IP_Antibody_Therapeutics_Landscape_Analysis.docx"
doc.save(docx_path)
print(f"Saved DOCX: {docx_path}")

# ============================================================================
# PDF BUILD (reportlab) - independent build, no LibreOffice conversion available
# ============================================================================
print("Building PDF ...")

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(name="TitleBig", fontSize=26, leading=30, alignment=1, textColor=colors.HexColor("#1F2A44"), spaceAfter=10))
styles.add(ParagraphStyle(name="SubtitleBig", fontSize=13, leading=16, alignment=1, textColor=colors.HexColor("#444444"), spaceAfter=6, fontName="Helvetica-Oblique"))
styles.add(ParagraphStyle(name="MetaCenter", fontSize=10, leading=14, alignment=1, textColor=colors.HexColor("#666666")))
styles.add(ParagraphStyle(name="H1", fontSize=17, leading=21, spaceBefore=16, spaceAfter=8, textColor=colors.HexColor("#1F2A44")))
styles.add(ParagraphStyle(name="H2", fontSize=13, leading=17, spaceBefore=12, spaceAfter=6, textColor=colors.HexColor("#2563EB")))
styles.add(ParagraphStyle(name="Body9", fontSize=9.5, leading=13, spaceAfter=6))
styles.add(ParagraphStyle(name="Caption", fontSize=8, leading=10, textColor=colors.HexColor("#666666"), spaceAfter=10, alignment=1))
styles.add(ParagraphStyle(name="Bullet9", fontSize=9.5, leading=13, leftIndent=14, bulletIndent=2, spaceAfter=4))

story = []
story.append(Spacer(1, 1.6 * inch))
story.append(Paragraph(REPORT_TITLE, styles["TitleBig"]))
story.append(Paragraph(SUBTITLE, styles["SubtitleBig"]))
story.append(Spacer(1, 0.3 * inch))
story.append(Paragraph(
    f"Report generated: {TODAY}<br/>Source notebook: ip_lanscape.ipynb<br/>"
    f"Source data: FINAL_master_plus_mab_addendum.csv ({TOTAL_ROWS:,} records)",
    styles["MetaCenter"],
))
story.append(PageBreak())

story.append(Paragraph("Executive Summary", styles["H1"]))
story.append(Paragraph(EXEC_SUMMARY, styles["Body9"]))
story.append(Paragraph("TL;DR", styles["H1"]))
for b in TLDR_BULLETS:
    story.append(Paragraph(f"&bull; {b}", styles["Bullet9"]))
story.append(Paragraph("Most Crowded Targets", styles["H2"]))
story.append(Paragraph(f"<b>Most crowded targets:</b> {CROWDED_LIST_STR}", styles["Body9"]))
story.append(Paragraph("White Space Targets Surfaced", styles["H2"]))
story.append(Paragraph(
    "This analysis has surfaced clear white space targets - see the Emerging Whitespace Deep-Dive and White Space "
    f"Targets per Solid Tumor Indication sections below. Standout unclaimed targets: <b>{STANDOUT_WHITESPACE}</b>.",
    styles["Body9"],
))
story.append(PageBreak())


def rl_df_table(df, max_rows=10, col_widths=None):
    shown = df.head(max_rows)
    data = [list(shown.columns)] + [
        ["" if pd_isna(v) else str(v) for v in row] for row in shown.itertuples(index=False)
    ]
    tbl = RLTable(data, hAlign="LEFT", colWidths=col_widths, repeatRows=1)
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2563EB")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, -1), 7),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CCCCCC")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F3F4F6")]),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    return tbl


def rl_add_table_section(story, entry, note=""):
    story.append(Paragraph(entry["title"], styles["H2"]))
    if note:
        story.append(Paragraph(note, styles["Body9"]))
    story.append(rl_df_table(entry["df"]))
    if entry["csv_path"] is not None:
        rel = f"data/{entry['csv_path'].name}"
        story.append(Paragraph(
            f"Full table ({len(entry['df']):,} rows): <link href='{rel}' color='blue'>{rel}</link>",
            styles["Caption"],
        ))
    story.append(Spacer(1, 10))


def rl_fig(path, width=6.3 * inch, caption=None):
    img = RLImage(str(path))
    ratio = img.imageHeight / float(img.imageWidth)
    img.drawWidth = width
    img.drawHeight = width * ratio
    story.append(img)
    if caption:
        story.append(Paragraph(caption, styles["Caption"]))


story.append(Paragraph("1. Dataset Overview & Methodology", styles["H1"]))
story.append(Paragraph(
    f"The underlying dataset contains {TOTAL_ROWS:,} patent records related to antibody therapeutics in oncology, "
    f"covering filing years {int(ipl.overview_df['min_year'].iloc[0])}-{int(ipl.overview_df['max_year'].iloc[0])}. "
    "All analytics are produced via DuckDB SQL queries executed directly against the source CSV, mirroring the "
    "queries documented in the Code Log of ip_lanscape.ipynb.",
    styles["Body9"],
))
rl_fig(ipl.fig01_path, caption="Figure 1. Overview dashboard: filing trend, modality mix, top targets, top assignees, indication whitespace and target x modality crowding matrix.")

story.append(Paragraph("2. Modality Distribution", styles["H2"]))
story.append(rl_df_table(ipl.modality_df, max_rows=10))
story.append(Spacer(1, 10))

for key in ["table04_filing_trend_by_year", "table01_top25_targets", "table03_top25_assignees",
            "table02_top25_indication_strings", "table05_claim_focus_top15", "table06_crowding_relevance_top15",
            "table07_parsed_indication_whitespace", "table08_target_modality_crowding_matrix_top15"]:
    entry = next(t for t in ipl.TABLE_REGISTRY if t["key"] == key)
    rl_add_table_section(story, entry)

story.append(PageBreak())
story.append(Paragraph("3. Target Whitespace Scoring & Multi-Angle View", styles["H1"]))
rl_add_table_section(story, next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table09_target_whitespace_scoring_top30"))
rl_fig(ipl.fig02_path, caption="Figure 2. Whitespace quadrant scatter, target x indication crowding heatmap, modality share and indication-resolution rate.")
for key in ["table10A_target_whitespace_summary_all", "table10B_target_modality_share_all",
            "table10C_target_indication_crowding_matrix", "table10D_whitespace_quadrant_classification"]:
    entry = next(t for t in ipl.TABLE_REGISTRY if t["key"] == key)
    rl_add_table_section(story, entry)

story.append(PageBreak())
story.append(Paragraph("4. Emerging Whitespace Deep-Dive", styles["H1"]))
story.append(Paragraph(
    f"{N_EMERGING} targets meet the Emerging Whitespace definition (total patents &lt; 40 and &gt;=30% filed since "
    f"2023). The table below shows the top 10 by volume; the full set is in the linked CSV.",
    styles["Body9"],
))
story.append(rl_df_table(EMERGING_TOP10[["target_name", "total_patents", "pct_recent", "top_assignee", "dominant_modality"]], max_rows=10))
entry = next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table11E_emerging_whitespace_deepdive")
rel = f"data/{entry['csv_path'].name}"
story.append(Paragraph(f"Full table ({len(entry['df']):,} rows): <link href='{rel}' color='blue'>{rel}</link>", styles["Caption"]))
rl_fig(ipl.fig03_path, caption="Figure 3. Emerging whitespace volume/momentum, filing trend, assignee concentration and dominant modality share.")
rl_add_table_section(story, next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table11F_emerging_whitespace_filing_trend"),
                      note="Year-by-year filing counts for every emerging whitespace target (2015-2025).")

story.append(PageBreak())
story.append(Paragraph("5. Target x Modality x Indication Whitespace Grid", styles["H1"]))
story.append(Paragraph(
    f"Across the top 10 targets, top 5 modalities and top 8 indications ({N_TOTAL_CELLS_G} combinations), "
    f"{N_ZERO_CELLS_G} combinations ({100*N_ZERO_CELLS_G/N_TOTAL_CELLS_G:.0f}%) have zero filed patents.",
    styles["Body9"],
))
rl_fig(ipl.fig04a_path, caption="Figure 4a. Small-multiple heatmaps of modality x indication patent counts per target (green/0 = white space, red = crowded).")
rl_fig(ipl.fig04b_path, width=5.5 * inch, caption="Figure 4b. Whitespace breadth (unclaimed modality x indication cells) per target.")
rl_add_table_section(story, next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table12G_target_modality_indication_grid"))

story.append(PageBreak())
story.append(Paragraph("6. White Space Targets Available per Solid Tumor Indication", styles["H1"]))
story.append(Paragraph(
    f"Considering the 20 most-established targets against the 12 leading solid tumor indications "
    f"({N_TOTAL_CELLS_H} combinations), {N_ZERO_CELLS_H} combinations ({100*N_ZERO_CELLS_H/N_TOTAL_CELLS_H:.0f}%) "
    f"are entirely unclaimed. Standout white space targets: <b>{STANDOUT_WHITESPACE}</b>.",
    styles["Body9"],
))
story.append(rl_df_table(WHITESPACE_INDICATIONS, max_rows=12))
rl_fig(ipl.fig05_path, caption="Figure 5. Solid tumor indication x target whitespace heatmap and whitespace breadth per indication.")
entry = next(t for t in ipl.TABLE_REGISTRY if t["key"] == "table13H_solid_tumor_indication_target_grid")
rel = f"data/{entry['csv_path'].name}"
story.append(Paragraph(f"Full grid ({len(entry['df']):,} rows): <link href='{rel}' color='blue'>{rel}</link>", styles["Caption"]))

story.append(PageBreak())
story.append(Paragraph("7. Key Findings & Recommendations", styles["H1"]))
for r in RECOMMENDATIONS:
    story.append(Paragraph(f"&bull; {r}", styles["Bullet9"]))

story.append(Paragraph("8. Data Quality, Caveats & Methodology Notes", styles["H1"]))
for n in DATA_QUALITY_NOTES:
    story.append(Paragraph(f"&bull; {n}", styles["Bullet9"]))

story.append(Paragraph("9. Making This Report Stand Out - Suggestions", styles["H1"]))
for s in STANDOUT_SUGGESTIONS:
    story.append(Paragraph(f"&bull; {s}", styles["Bullet9"]))

story.append(PageBreak())
story.append(Paragraph("Appendix: Linked Data Files", styles["H1"]))
story.append(Paragraph(
    "Every table above with more than 10 rows is exported in full to the data/ folder as an individual CSV. All "
    "tables are also consolidated into a single multi-sheet Excel workbook: "
    f"<link href='data/{ipl.xlsx_path.name}' color='blue'>data/{ipl.xlsx_path.name}</link>",
    styles["Body9"],
))
appendix_data = [["Table", "Rows", "File"]] + [
    [t["title"], str(len(t["df"])), f"data/{t['csv_path'].name}"] for t in linked_tables()
]
appendix_tbl = RLTable(appendix_data, hAlign="LEFT", colWidths=[2.6 * inch, 0.6 * inch, 3.1 * inch], repeatRows=1)
appendix_tbl.setStyle(TableStyle([
    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2563EB")),
    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
    ("FONTSIZE", (0, 0), (-1, -1), 7.5),
    ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CCCCCC")),
    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F3F4F6")]),
]))
story.append(appendix_tbl)

pdf_path = OUT / "IP_Antibody_Therapeutics_Landscape_Analysis.pdf"
pdf_doc = SimpleDocTemplate(str(pdf_path), pagesize=letter, topMargin=0.6 * inch, bottomMargin=0.6 * inch,
                            leftMargin=0.6 * inch, rightMargin=0.6 * inch, title=REPORT_TITLE)
pdf_doc.build(story)
print(f"Saved PDF: {pdf_path}")

# ============================================================================
# PPTX BUILD (python-pptx)
# ============================================================================
print("Building PPTX ...")

prs = Presentation()
prs.slide_width = PIn(13.333)
prs.slide_height = PIn(7.5)
BLANK = prs.slide_layouts[6]
DARK = PRGBColor(0x1F, 0x2A, 0x44)
ACCENT = PRGBColor(0x25, 0x63, 0xEB)
GREY = PRGBColor(0x44, 0x44, 0x44)


def add_title_box(slide, text, top=PIn(0.35), size=32, color=DARK, height=PIn(1.0)):
    box = slide.shapes.add_textbox(PIn(0.5), top, prs.slide_width - PIn(1.0), height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PPt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return box


def add_bullets(slide, bullets, top=PIn(1.5), left=PIn(0.7), width=None, height=None, size=16):
    width = width or (prs.slide_width - PIn(1.4))
    height = height or (prs.slide_height - top - PIn(0.4))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = PPt(size)
        p.font.color.rgb = GREY
        p.space_after = PPt(10)
    return box


def add_image_slide(title, image_path, caption=None):
    slide = prs.slides.add_slide(BLANK)
    add_title_box(slide, title, size=26, height=PIn(0.7))
    from PIL import Image as PILImage
    with PILImage.open(image_path) as im:
        w, h = im.size
    ratio = h / w
    max_w = PIn(12.2)
    max_h = PIn(5.9)
    draw_w = max_w
    draw_h = PIn(max_w.inches * ratio)
    if draw_h > max_h:
        draw_h = max_h
        draw_w = PIn(max_h.inches / ratio)
    left = (prs.slide_width - draw_w) / 2
    slide.shapes.add_picture(str(image_path), left, PIn(1.15), width=draw_w, height=draw_h)
    if caption:
        cap = slide.shapes.add_textbox(PIn(0.5), prs.slide_height - PIn(0.55), prs.slide_width - PIn(1.0), PIn(0.4))
        p = cap.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = PPt(11)
        p.font.italic = True
        p.font.color.rgb = GREY
        p.alignment = PP_ALIGN.CENTER
    return slide


# Slide 1: Title
slide = prs.slides.add_slide(BLANK)
box = slide.shapes.add_textbox(PIn(0.8), PIn(2.6), prs.slide_width - PIn(1.6), PIn(1.6))
p = box.text_frame.paragraphs[0]
p.text = REPORT_TITLE
p.font.size = PPt(40)
p.font.bold = True
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER
box.text_frame.word_wrap = True
sub = slide.shapes.add_textbox(PIn(0.8), PIn(4.1), prs.slide_width - PIn(1.6), PIn(1.0))
p = sub.text_frame.paragraphs[0]
p.text = SUBTITLE
p.font.size = PPt(18)
p.font.italic = True
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER
sub.text_frame.word_wrap = True
meta = slide.shapes.add_textbox(PIn(0.8), PIn(6.4), prs.slide_width - PIn(1.6), PIn(0.8))
p = meta.text_frame.paragraphs[0]
p.text = f"{TODAY}  |  Source: ip_lanscape.ipynb  |  {TOTAL_ROWS:,} patent records"
p.font.size = PPt(13)
p.font.color.rgb = GREY
p.alignment = PP_ALIGN.CENTER

# Slide 2: Executive Summary
slide = prs.slides.add_slide(BLANK)
add_title_box(slide, "Executive Summary")
add_bullets(slide, [EXEC_SUMMARY], size=14)

# Slide 3: TL;DR
slide = prs.slides.add_slide(BLANK)
add_title_box(slide, "TL;DR")
add_bullets(slide, TLDR_BULLETS, size=15)

# Slide 4: Most Crowded Targets
slide = prs.slides.add_slide(BLANK)
add_title_box(slide, "Most Crowded Targets")
crowded_bullets = [f"{r.target_name}: {int(r.n)} patents" for r in CROWDED_TOP10.itertuples()]
add_bullets(slide, crowded_bullets, size=16)

# Slide 5: White Space Findings
slide = prs.slides.add_slide(BLANK)
add_title_box(slide, "White Space Targets Surfaced")
whitespace_bullets = [
    f"{N_EMERGING} targets classified as Emerging Whitespace (low volume, high recent momentum)",
    "Top emerging targets by volume: " + ", ".join(EMERGING_TOP10["target_name"].head(6).tolist()),
    f"{N_ZERO_CELLS_H}/{N_TOTAL_CELLS_H} target x solid-tumor-indication combinations are entirely unclaimed (top 20 targets x top 12 indications)",
    f"Standout unclaimed targets across indications: {STANDOUT_WHITESPACE}",
]
add_bullets(slide, whitespace_bullets, size=16)

# Chart slides
add_image_slide("Overview Dashboard", ipl.fig01_path,
                "Filing trend, modality mix, top targets, top assignees, indication whitespace, target x modality crowding matrix")
add_image_slide("Target Whitespace: Multi-Angle View", ipl.fig02_path,
                "Whitespace quadrant, target x indication heatmap, modality share, indication-resolution rate")
add_image_slide("Emerging Whitespace Deep-Dive", ipl.fig03_path,
                "Volume/momentum, filing trend, assignee concentration, dominant modality share")
add_image_slide("Target x Modality x Indication Whitespace Grid", ipl.fig04a_path,
                "Small-multiple heatmaps per target (green/0 = white space, red = crowded)")
add_image_slide("Whitespace Breadth by Target", ipl.fig04b_path,
                "Count of unclaimed modality x indication combinations per target")
add_image_slide("Solid Tumor Indication White Space", ipl.fig05_path,
                "Solid tumor indication x target whitespace heatmap and breadth per indication")

# Recommendations slide
slide = prs.slides.add_slide(BLANK)
add_title_box(slide, "Key Findings & Recommendations")
add_bullets(slide, RECOMMENDATIONS, size=15)

# Stand-out suggestions slide
slide = prs.slides.add_slide(BLANK)
add_title_box(slide, "Making This Report Stand Out")
add_bullets(slide, STANDOUT_SUGGESTIONS, size=14)

# Appendix slide
slide = prs.slides.add_slide(BLANK)
add_title_box(slide, "Appendix: Full Data Exports")
appendix_bullets = [f"{ipl.xlsx_path.name} - combined workbook, all {len(ipl.TABLE_REGISTRY)} tables"] + [
    f"{t['csv_path'].name} ({len(t['df']):,} rows)" for t in linked_tables()
]
add_bullets(slide, appendix_bullets, size=11)

pptx_path = OUT / "IP_Antibody_Therapeutics_Landscape_Analysis.pptx"
prs.save(pptx_path)
print(f"Saved PPTX: {pptx_path}")

print("\nAll deliverables built successfully in:", OUT)
