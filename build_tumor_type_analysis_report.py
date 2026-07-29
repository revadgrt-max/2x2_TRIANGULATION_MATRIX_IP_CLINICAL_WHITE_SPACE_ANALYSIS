"""
build_tumor_type_analysis_report.py
=====================================
Standalone report covering ONLY the Part C "Tumor-Type Deep Dive" section
(Entries 19-23) of the already-executed clinical_ip_comparison_final_version1.ipynb
-- the indication-level re-cut of the target-level Part B triangulation
(validated/patent-thin, patent-crowded/unproven, timing lag, and the CAR-T
blind spot, each broken down by solid-tumor indication).

Parses the notebook directly (no recomputation), so every number/table/chart
here matches the notebook exactly.

Produces, under output/clinical_ip_comparison_version1/:
    Tumor_Type_Analysis_Report.pdf
    Tumor_Type_Analysis_Report.docx
    Tumor_Type_Analysis_Deck.pptx
"""
import base64
import io
import json
import re
from datetime import date
from pathlib import Path

import pandas as pd

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table as RLTable,
    TableStyle, PageBreak,
)
from PIL import Image as PILImage

from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

BASE = Path("/Users/revathisekar/Documents/vibe analytics trial 2")
NOTEBOOK_SRC = BASE / "input" / "clinical_ip_comparison_final_version1.ipynb"
OUT = BASE / "output" / "clinical_ip_comparison_version1"
DATA = OUT / "data"
FIG = OUT / "figures"
DATA.mkdir(parents=True, exist_ok=True)
FIG.mkdir(parents=True, exist_ok=True)

TODAY = date(2026, 7, 27).strftime("%B %d, %Y")
REPORT_TITLE = "Tumor-Type Deep Dive -- Clinical x IP Triangulation, by Solid-Tumor Indication"
SUBTITLE = "Indication-level re-cut of the target-level triangulation: validated/patent-thin, FTO risk, timing lag, and the CAR-T blind spot"


def pd_isna(v):
    try:
        return pd.isna(v)
    except (TypeError, ValueError):
        return False


print("Loading and parsing notebook ...")
nb = json.loads(NOTEBOOK_SRC.read_text())


def cell_source(cell):
    src = cell.get("source", "")
    return "".join(src) if isinstance(src, list) else src


ENTRY_RE = re.compile(r"^##\s+Entry\s+(\d+)\s*(?:\[([^\]]*)\])?\s*:\s*(.+)$")

all_entries = []
current = None
_fig_counter = 0

for cell in nb["cells"]:
    ctype = cell.get("cell_type")
    src = cell_source(cell)
    if ctype == "markdown":
        stripped = src.strip()
        lines = stripped.split("\n")
        first_line = lines[0].strip() if lines else ""
        m = ENTRY_RE.match(first_line)
        if m:
            if current is not None:
                all_entries.append(current)
            number = int(m.group(1))
            tag = (m.group(2) or "").strip()
            title = m.group(3).strip()
            desc = "\n".join(lines[1:]).strip()
            current = {"number": number, "tag": tag, "title": title, "description": desc, "items": [], "code": ""}
        continue
    elif ctype == "code":
        if current is None:
            continue
        current["code"] = cell_source(cell)
        for out in cell.get("outputs", []):
            otype = out.get("output_type")
            data = out.get("data", {})
            meta = out.get("metadata", {}) or {}
            if otype in ("display_data", "execute_result"):
                if "image/png" in data:
                    b64 = data["image/png"]
                    if isinstance(b64, list):
                        b64 = "".join(b64)
                    _fig_counter += 1
                    fname = f"entry{current['number']:02d}_fig{_fig_counter:02d}.png"
                    fpath = FIG / fname
                    if not fpath.exists():
                        fpath.write_bytes(base64.b64decode(b64))
                    current["items"].append({"type": "image", "path": fpath})
                    continue
                html = data.get("text/html")
                if html:
                    html = "".join(html) if isinstance(html, list) else html
                    try:
                        dfs = pd.read_html(io.StringIO(html))
                        tdf = dfs[0] if dfs else None
                    except Exception:
                        tdf = None
                    if tdf is not None and len(tdf.columns) > 0:
                        tname = meta.get("name") or f"entry{current['number']:02d}_table"
                        safe = re.sub(r"[^A-Za-z0-9_]+", "_", tname)
                        csv_name = f"{safe}.csv"
                        csv_path = DATA / csv_name
                        if not csv_path.exists():
                            tdf.to_csv(csv_path, index=False)
                        current["items"].append({"type": "table", "df": tdf, "csv_name": csv_name, "title": tname})
                        continue
                txt = data.get("text/plain")
                if txt:
                    txt = "".join(txt) if isinstance(txt, list) else txt
                    if txt.strip():
                        current["items"].append({"type": "text", "text": txt.strip()})
                        continue
            elif otype == "stream":
                txt_raw = out.get("text", "")
                txt = "".join(txt_raw) if isinstance(txt_raw, list) else txt_raw
                if txt.strip():
                    current["items"].append({"type": "text", "text": txt.strip()})

if current is not None:
    all_entries.append(current)

# Only keep Part C (tumor-type deep dive) entries.
entries = [e for e in all_entries if e["number"] >= 19]
print(f"Parsed {len(all_entries)} total entries; keeping {len(entries)} Part-C tumor-type entries (19-{entries[-1]['number']}).")

e23 = next((e for e in entries if e["number"] == 23), None)
exec_summary = next((it["text"] for it in (e23["items"] if e23 else []) if it["type"] == "text" and "TUMOR-TYPE ANALYSIS" in it["text"]), "")

TLDR_BULLETS = [
    "Part B of the main triangulation (Entries 10-16) was computed at the TARGET level, pooling every "
    "indication together. This deep dive re-cuts the two highest-value business questions -- 'file now' "
    "and 'FTO risk' -- plus timing lag and the CAR-T blind spot, at the (target, solid-tumor-indication) level.",
    "Uses the same canonical solid-tumor taxonomy applied consistently across the IP `indications` free text "
    "and the clinical `ctgov_conditions`/`m_conditions` fields (Entry 2/9 of the main notebook), so patent and "
    "trial volume are directly comparable indication-by-indication.",
    "Entry 19: (target, indication) pairs with a recorded clinical success AND <=5 patents in that same "
    "indication -- the sharpest indication-specific 'file now' case.",
    "Entry 20: (target, indication) pairs with >=5 patents but zero clinical success in that same indication -- "
    "the sharpest indication-specific freedom-to-operate / litigation-exposure view.",
    "Entry 21: filing-to-first-trial lag recomputed per indication -- flags cases where clinical trials in one "
    "specific indication started before any patent was filed, even if the target overall looks IP-led.",
    "Entry 22: CAR-T clinical activity cross-referenced against antibody-modality IP coverage per indication -- "
    "deliberately left unfiltered by tumor category since CAR-T's real indication mix is mostly hematologic, "
    "which the solid-tumor-only entries elsewhere in the notebook exclude.",
]

# ============================================================================
# PDF
# ============================================================================
print("Building PDF ...")
styles = getSampleStyleSheet()
styles.add(ParagraphStyle(name="TitleBig", fontSize=20, leading=24, alignment=1, textColor=colors.HexColor("#1F2A44"), spaceAfter=10))
styles.add(ParagraphStyle(name="SubtitleBig", fontSize=11.5, leading=15, alignment=1, textColor=colors.HexColor("#444444"), spaceAfter=6, fontName="Helvetica-Oblique"))
styles.add(ParagraphStyle(name="MetaCenter", fontSize=9.5, leading=13, alignment=1, textColor=colors.HexColor("#666666")))
styles.add(ParagraphStyle(name="H1", fontSize=15, leading=19, spaceBefore=14, spaceAfter=8, textColor=colors.HexColor("#1F2A44")))
styles.add(ParagraphStyle(name="Body9", fontSize=9, leading=12.5, spaceAfter=6))
styles.add(ParagraphStyle(name="Note9", fontSize=9, leading=12.5, spaceAfter=6, textColor=colors.HexColor("#333333"), fontName="Helvetica-Oblique"))
styles.add(ParagraphStyle(name="Mono8", fontSize=7.5, leading=10.5, spaceAfter=6, fontName="Courier"))
styles.add(ParagraphStyle(name="Bullet9", fontSize=9, leading=12.5, leftIndent=14, bulletIndent=2, spaceAfter=4))
styles.add(ParagraphStyle(name="TOCLink", fontSize=10, leading=15, textColor=colors.HexColor("#1155CC")))


def esc(text):
    return str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def esc_pre(text):
    return esc(text).replace("\n", "<br/>")


def pdf_df_table(df, max_rows=20, font_size=6.0):
    shown = df.head(max_rows)
    data = [[esc(c) for c in shown.columns]] + [
        ["" if pd_isna(v) else esc(v) for v in row] for row in shown.itertuples(index=False)
    ]
    ncols = len(shown.columns)
    avail_width = 7.3 * inch
    col_width = avail_width / max(ncols, 1)
    t = RLTable(data, colWidths=[col_width] * ncols, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F2A44")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, -1), font_size),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CCCCCC")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F6FA")]),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 3),
        ("RIGHTPADDING", (0, 0), (-1, -1), 3),
    ]))
    return t


def pdf_render_items(story, items, max_rows=25):
    for it in items:
        if it["type"] == "text":
            style = styles["Mono8"] if len(it["text"]) > 300 else styles["Note9"]
            story.append(Paragraph(esc_pre(it["text"]), style))
            story.append(Spacer(1, 4))
        elif it["type"] == "table":
            story.append(pdf_df_table(it["df"], max_rows=max_rows))
            rel = f"data/{it['csv_name']}"
            if len(it["df"]) > max_rows:
                story.append(Paragraph(
                    f"Showing first {max_rows} of {len(it['df']):,} rows. "
                    f"<link href=\"{rel}\">Click to open full CSV ({esc(it['csv_name'])})</link>", styles["Body9"]))
            else:
                story.append(Paragraph(
                    f"All {len(it['df']):,} rows shown above. <link href=\"{rel}\">Click to open CSV ({esc(it['csv_name'])})</link>",
                    styles["Body9"]))
            story.append(Spacer(1, 8))
        elif it["type"] == "image":
            with PILImage.open(it["path"]) as im:
                w, h = im.size
            max_w = 6.6 * inch
            ratio = h / w
            draw_w = max_w
            draw_h = max_w * ratio
            max_h = 7.3 * inch
            if draw_h > max_h:
                draw_h = max_h
                draw_w = max_h / ratio
            story.append(RLImage(str(it["path"]), width=draw_w, height=draw_h))
            story.append(Spacer(1, 8))


def entry_label(e):
    tag = f" [{e['tag']}]" if e["tag"] else ""
    return f"Entry {e['number']}{tag}: {e['title']}"


toc_entries = [("exec_summary", "Executive Summary"), ("tldr", "TL;DR")]
toc_entries += [(f"sec_entry_{e['number']}", entry_label(e)) for e in entries]
toc_entries += [("appendix", "Appendix: Data Files Index")]

story = []
story.append(Spacer(1, 1.2 * inch))
story.append(Paragraph(esc(REPORT_TITLE), styles["TitleBig"]))
story.append(Paragraph(esc(SUBTITLE), styles["SubtitleBig"]))
story.append(Spacer(1, 0.25 * inch))
story.append(Paragraph(
    f"Report generated: {TODAY}<br/>Source notebook: clinical_ip_comparison_final_version1.ipynb "
    f"(Entries 19-{entries[-1]['number']})<br/>Companion to: Clinical_IP_Comparison_Report.pdf (full triangulation)",
    styles["MetaCenter"]))
story.append(PageBreak())

story.append(Paragraph('<a name="toc"/>Table of Contents', styles["H1"]))
for anchor, label in toc_entries:
    story.append(Paragraph(f'<link href="#{anchor}">{esc(label)}</link>', styles["TOCLink"]))
story.append(PageBreak())

story.append(Paragraph('<a name="exec_summary"/>Executive Summary', styles["H1"]))
if exec_summary:
    story.append(Paragraph(esc_pre(exec_summary), styles["Mono8"]))
story.append(PageBreak())

story.append(Paragraph('<a name="tldr"/>TL;DR', styles["H1"]))
for b in TLDR_BULLETS:
    story.append(Paragraph(f"&bull; {esc(b)}", styles["Bullet9"]))
story.append(PageBreak())

for e in entries:
    story.append(Paragraph(f'<a name="sec_entry_{e["number"]}"/>{esc(entry_label(e))}', styles["H1"]))
    if e["description"]:
        story.append(Paragraph(esc(e["description"]), styles["Body9"]))
    pdf_render_items(story, e["items"])
    story.append(PageBreak())

story.append(Paragraph('<a name="appendix"/>Appendix: Data Files Index', styles["H1"]))
story.append(Paragraph(esc(
    "Every table in this report has its full underlying data deposited as a CSV file in the data/ folder "
    "next to this report (shared with the main Clinical_IP_Comparison_Report.pdf)."), styles["Body9"]))
appendix_rows = [["Entry", "Rows", "File"]]
for e in entries:
    for it in e["items"]:
        if it["type"] == "table":
            appendix_rows.append([entry_label(e)[:45], f"{len(it['df']):,}", it["csv_name"]])
t = RLTable(appendix_rows, colWidths=[3.4 * inch, 0.7 * inch, 3.2 * inch], repeatRows=1)
t.setStyle(TableStyle([
    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F2A44")),
    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
    ("FONTSIZE", (0, 0), (-1, -1), 7),
    ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CCCCCC")),
    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F6FA")]),
]))
story.append(t)

pdf_path = OUT / "Tumor_Type_Analysis_Report.pdf"
pdf_doc = SimpleDocTemplate(str(pdf_path), pagesize=letter, topMargin=0.6 * inch, bottomMargin=0.6 * inch,
                            leftMargin=0.6 * inch, rightMargin=0.6 * inch, title=REPORT_TITLE)
pdf_doc.build(story)
print(f"Saved PDF: {pdf_path}")

# ============================================================================
# DOCX
# ============================================================================
print("Building DOCX ...")


def add_hyperlink(paragraph, text, target_path):
    part = paragraph.part
    r_id = part.relate_to(
        target_path, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    new_run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "1155CC")
    rPr.append(color)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rPr.append(u)
    new_run.append(rPr)
    t = OxmlElement("w:t")
    t.text = text
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)
    return hyperlink


def docx_add_df_table(doc, df, max_rows=25):
    shown = df.head(max_rows)
    table = doc.add_table(rows=1, cols=len(shown.columns))
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr_cells = table.rows[0].cells
    for i, col in enumerate(shown.columns):
        hdr_cells[i].text = str(col)
        for p in hdr_cells[i].paragraphs:
            for r in p.runs:
                r.bold = True
    for _, row in shown.iterrows():
        cells = table.add_row().cells
        for i, val in enumerate(row):
            cells[i].text = "" if pd_isna(val) else str(val)
    return table


def docx_render_items(doc, items, max_rows=25):
    for it in items:
        if it["type"] == "text":
            doc.add_paragraph(it["text"])
        elif it["type"] == "table":
            docx_add_df_table(doc, it["df"], max_rows=max_rows)
            p = doc.add_paragraph()
            rel = f"data/{it['csv_name']}"
            if len(it["df"]) > max_rows:
                p.add_run(f"Showing first {max_rows} of {len(it['df']):,} rows. Full table: ").italic = True
            else:
                p.add_run(f"All {len(it['df']):,} rows shown above. Full table: ").italic = True
            add_hyperlink(p, rel, rel)
            doc.add_paragraph()
        elif it["type"] == "image":
            doc.add_picture(str(it["path"]), width=Inches(6.2))
            doc.add_paragraph()


doc = Document()
title_p = doc.add_paragraph()
title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = title_p.add_run(REPORT_TITLE)
run.bold = True
run.font.size = Pt(24)
run.font.color.rgb = RGBColor(0x1F, 0x2A, 0x44)

sub_p = doc.add_paragraph()
sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = sub_p.add_run(SUBTITLE)
run.font.size = Pt(12.5)
run.italic = True

meta_p = doc.add_paragraph()
meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
meta_p.add_run(f"Report generated: {TODAY}\nSource notebook: clinical_ip_comparison_final_version1.ipynb "
               f"(Entries 19-{entries[-1]['number']})\nCompanion to: Clinical_IP_Comparison_Report.docx")
doc.add_page_break()

doc.add_heading("Executive Summary", level=1)
doc.add_paragraph(exec_summary)
doc.add_heading("TL;DR", level=1)
for b in TLDR_BULLETS:
    doc.add_paragraph(b, style="List Bullet")
doc.add_page_break()

for e in entries:
    doc.add_heading(entry_label(e), level=1)
    if e["description"]:
        doc.add_paragraph(e["description"])
    docx_render_items(doc, e["items"], max_rows=25)
    doc.add_page_break()

doc.add_heading("Appendix: Data Files Index", level=1)
doc.add_paragraph("Every table in this report has its full underlying data deposited as a CSV file in the "
                   "data/ folder next to this report (shared with the main report).")
app_table = doc.add_table(rows=1, cols=3)
app_table.style = "Light Grid Accent 1"
for i, h in enumerate(["Entry", "Rows", "File"]):
    app_table.rows[0].cells[i].text = h
for e in entries:
    for it in e["items"]:
        if it["type"] == "table":
            cells = app_table.add_row().cells
            cells[0].text = entry_label(e)[:60]
            cells[1].text = f"{len(it['df']):,}"
            cells[2].text = it["csv_name"]

docx_path = OUT / "Tumor_Type_Analysis_Report.docx"
doc.save(docx_path)
print(f"Saved DOCX: {docx_path}")

# ============================================================================
# PPTX -- 16:9 deck with a clickable agenda
# ============================================================================
print("Building PPTX ...")

prs = Presentation()
prs.slide_width = PIn(13.333)
prs.slide_height = PIn(7.5)
BLANK = prs.slide_layouts[6]
DARK = PRGBColor(0x1F, 0x2A, 0x44)
ACCENT = PRGBColor(0x25, 0x63, 0xEB)
GREY = PRGBColor(0x44, 0x44, 0x44)


def add_title_box(slide, text, top=PIn(0.35), size=24, color=DARK, height=PIn(1.0)):
    box = slide.shapes.add_textbox(PIn(0.5), top, prs.slide_width - PIn(1.0), height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PPt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return box


def add_bullets(slide, bullets, top=PIn(1.5), left=PIn(0.7), width=None, height=None, size=15):
    width = width or (prs.slide_width - PIn(1.4))
    height = height or (prs.slide_height - top - PIn(0.4))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022 {b}"
        p.font.size = PPt(size)
        p.font.color.rgb = GREY
        p.space_after = PPt(10)
    return box


def add_image_slide(title, image_path, caption=None):
    slide = prs.slides.add_slide(BLANK)
    add_title_box(slide, title, size=19, height=PIn(0.7))
    with PILImage.open(image_path) as im:
        w, h = im.size
    ratio = h / w
    max_w = PIn(12.2)
    max_h = PIn(5.9)
    draw_w = max_w
    draw_h = PIn(max_w.inches * ratio)
    if draw_h > max_h:
        draw_h = max_h
        draw_w = PIn(max_h.inches / ratio)
    left = (prs.slide_width - draw_w) / 2
    slide.shapes.add_picture(str(image_path), left, PIn(1.15), width=draw_w, height=draw_h)
    if caption:
        cap = slide.shapes.add_textbox(PIn(0.5), prs.slide_height - PIn(0.55), prs.slide_width - PIn(1.0), PIn(0.4))
        p = cap.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = PPt(11)
        p.font.italic = True
        p.font.color.rgb = GREY
        p.alignment = PP_ALIGN.CENTER
    return slide


def add_table_slide(title, df, max_rows=8):
    slide = prs.slides.add_slide(BLANK)
    add_title_box(slide, title, size=19, height=PIn(0.7))
    shown = df.head(max_rows)
    rows, cols = shown.shape[0] + 1, shown.shape[1]
    left, top, width, height = PIn(0.6), PIn(1.3), prs.slide_width - PIn(1.2), PIn(5.6)
    gtable = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for j, col in enumerate(shown.columns):
        gtable.cell(0, j).text = str(col)
    for i, (_, row) in enumerate(shown.iterrows(), start=1):
        for j, val in enumerate(row):
            gtable.cell(i, j).text = "" if pd_isna(val) else str(val)
    return slide


def add_text_slide(title, text, size=13):
    slide = prs.slides.add_slide(BLANK)
    add_title_box(slide, title, size=22)
    lines = [ln for ln in text.split("\n") if ln.strip()][:22]
    add_bullets(slide, lines, size=size, height=PIn(5.8))
    return slide


slide = prs.slides.add_slide(BLANK)
box = slide.shapes.add_textbox(PIn(0.8), PIn(2.0), prs.slide_width - PIn(1.6), PIn(2.0))
p = box.text_frame.paragraphs[0]
p.text = REPORT_TITLE
p.font.size = PPt(30)
p.font.bold = True
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER
box.text_frame.word_wrap = True
sub_box = slide.shapes.add_textbox(PIn(0.8), PIn(4.0), prs.slide_width - PIn(1.6), PIn(1.4))
p2 = sub_box.text_frame.paragraphs[0]
p2.text = SUBTITLE
p2.font.size = PPt(15)
p2.font.italic = True
p2.font.color.rgb = GREY
p2.alignment = PP_ALIGN.CENTER
sub_box.text_frame.word_wrap = True
meta_box = slide.shapes.add_textbox(PIn(0.8), PIn(5.5), prs.slide_width - PIn(1.6), PIn(0.8))
p3 = meta_box.text_frame.paragraphs[0]
p3.text = f"Generated {TODAY}  |  clinical_ip_comparison_final_version1.ipynb (Entries 19-{entries[-1]['number']})"
p3.font.size = PPt(13)
p3.font.color.rgb = GREY
p3.alignment = PP_ALIGN.CENTER

add_text_slide("TL;DR", "\n".join(TLDR_BULLETS), size=14)

if exec_summary:
    lines = [ln for ln in exec_summary.split("\n") if ln.strip()]
    chunk_size = 16
    for i in range(0, len(lines), chunk_size):
        chunk = lines[i:i + chunk_size]
        title = "Executive Summary" if i == 0 else "Executive Summary (cont.)"
        add_text_slide(title, "\n".join(chunk), size=13)

agenda_slide = prs.slides.add_slide(BLANK)
add_title_box(agenda_slide, "Agenda", size=28)
agenda_top = PIn(1.4)
agenda_left_col = [PIn(0.6), PIn(6.9)]
col_width = PIn(6.0)
row_height = PIn(0.34)
n_per_col = 22

entry_first_slide = {}
for e in entries:
    images = [it for it in e["items"] if it["type"] == "image"]
    tables = [it for it in e["items"] if it["type"] == "table"]
    title = entry_label(e)
    made_slides = []
    if images:
        for k, img in enumerate(images):
            cap = title if k == 0 else f"{title} (cont.)"
            made_slides.append(add_image_slide(title, img["path"], caption=e["title"]))
    elif tables:
        made_slides.append(add_table_slide(title, tables[0]["df"], max_rows=8))
    else:
        notes = [it["text"] for it in e["items"] if it["type"] == "text"]
        text_blob = "\n".join(notes) if notes else "See full notebook / PDF report for details."
        made_slides.append(add_text_slide(title, text_blob, size=13))
    if made_slides:
        entry_first_slide[e["number"]] = made_slides[0]

for idx, e in enumerate(entries):
    col = idx // n_per_col
    row = idx % n_per_col
    left = agenda_left_col[min(col, len(agenda_left_col) - 1)]
    top = agenda_top + PIn(row_height.inches * row)
    box = agenda_slide.shapes.add_textbox(left, top, col_width, row_height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = entry_label(e)
    p.font.size = PPt(11)
    p.font.color.rgb = ACCENT
    target = entry_first_slide.get(e["number"])
    if target is not None:
        box.click_action.target_slide = target

pptx_path = OUT / "Tumor_Type_Analysis_Deck.pptx"
prs.save(pptx_path)
print(f"Saved PPTX: {pptx_path}")

print("\nDone.")
print(f"PDF:  {pdf_path}")
print(f"DOCX: {docx_path}")
print(f"PPTX: {pptx_path}")

